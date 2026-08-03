import datetime
import os.path
import time
import base64
import traceback
from email.mime.text import MIMEText
from email.message import EmailMessage # PROVEN LOGIC: Used for embedding images
import io # For GDrive downloads # For GDrive downloads
import re
import markdown 
import json
import fitz
from google import genai
from google.genai import types
import enum
import pandas as pd

from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from dotenv import load_dotenv
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaIoBaseDownload
from dotenv import load_dotenv
from pydantic import BaseModel, ValidationError
from data_config import sheet_masters, hierarchy, column_index

# For parsing Office documents if downloaded from Drive
from pptx import Presentation
from pptx.util import Inches
import openpyxl

# Additional helper imports for handling media uploads and memory streams in testing mode
from googleapiclient.http import MediaIoBaseUpload
import io

# Load the .env file
env = os.getenv("ENV", "dev")
env_file = f".env.{env}"

if not os.getenv("GITHUB_ACTIONS"):
    load_dotenv(env_file)

if not os.getenv("GITHUB_ACTIONS"):
    load_dotenv(env_file)
# --- Configuration ---
# For Google Workspace APIs (Calendar, Gmail, Drive)
SCOPES = [
    'https://www.googleapis.com/auth/calendar',
    'https://www.googleapis.com/auth/gmail.send',
    'https://www.googleapis.com/auth/drive.readonly',
    'https://www.googleapis.com/auth/spreadsheets',
    'https://www.googleapis.com/auth/drive',
    'https://www.googleapis.com/auth/documents'
]
CREDENTIALS_FILE = os.getenv("GOOGLE_CREDENTIALS_FILE", "credentials.json") # Downloaded from GCP
TOKEN_FILE_PREFIX = 'token_brandvmeet' # Will generate token_brandvmeet_calendar.json etc.

import requests # Serper.dev

# Google Drive Folder ID containing NBH data
NBH_GDRIVE_FOLDER_ID = os.getenv("NBH_GDRIVE_FOLDER_ID") # Set env var or replace placeholder

# For Gemini API
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY") 

# Load and sanitize Serper API Key automatically to strip trailing newlines, spaces, or accidental quotes
RAW_SERPER_KEY = os.getenv("SERPER_API_KEY", "")
if RAW_SERPER_KEY:
    SERPER_API_KEY = str(RAW_SERPER_KEY).strip().replace('"', '').replace("'", "")
else:
    SERPER_API_KEY = None

# =====================================================================
# SANITIZED SERPER.DEV SEARCH API CALL
# =====================================================================
def execute_serper_search_api(query, num_results=5):
    """
    Executes a Google search via Serper.dev pointing to India region.
    Cleans incoming inputs and safely logs request errors to help debug authentication issues.
    """
    if not SERPER_API_KEY:
        print("  ⚠️ [SERPER] API Key is missing or empty. Skipping search.")
        return "Search results unavailable: API Key missing."
        
    url = "https://google.serper.dev/search"
    payload = json.dumps({
        "q": query,
        "num": num_results,
        "gl": "in" # Constrain searches to India region
    })
    
    headers = {
        'X-API-KEY': SERPER_API_KEY,
        'Content-Type': 'application/json'
    }
    
    try:
        response = requests.post(url, headers=headers, data=payload, timeout=10)
        
        if response.status_code == 200:
            organic_results = response.json().get("organic", [])
            if not organic_results:
                return "No search results found on Google."
                
            formatted = []
            for i, res in enumerate(organic_results, 1):
                title = res.get("title", "No Title")
                link = res.get("link", "No Link")
                snippet = res.get("snippet", "")
                formatted.append(f"[{i}] {title}\nURL: {link}\nSnippet: {snippet}\n")
            return "\n".join(formatted)
            
        else:
            # Descriptive log output to immediately identify key issues in the action console
            print(f"  ⚠️ [SERPER] API Error. HTTP Status: {response.status_code} | Details: {response.text}")
            return f"Search failed. Code: {response.status_code}"
            
    except Exception as e:
        print(f"  ⚠️ [SERPER] Network or connection error: {e}")
        return f"Search network error: {e}"


AGENT_EMAIL = "brand.vmeet@nobroker.in" # Email of the agent account
ADMIN_EMAIL_FOR_NOTIFICATIONS = "ajay.saini@nobroker.in" # REPLACE with your actual email
leadership_emails = ["sristi.agarwal@nobroker.in", "rohit.c@nobroker.in"] # Add the second email

EXCLUDED_NBH_PSEUDO_NAMES_FOR_FOLLOWUP = {
    AGENT_EMAIL.lower().split('@')[0],
    "pia.brand","pia","nbh.meeting"
}

PROCESSED_EVENTS_FILE = 'processed_event_ids.txt' # Simple file-based tracking for local runs
# --- Feature Toggles ---
ENABLE_IMAGE_GENERATION = True  # Set to False to completely skip image logic



# Cache file for inferred industries
# INFERRED_INDUSTRIES_CACHE_FILE = 'inferred_industries_cache.json' # Cache for brand industry inference



# Specific File Names (you might make these configurable or discover them)
FILE_NAME_PITCH_DECK_PDF = "NBH Monetization Pitch Deck.pdf" # From your image
# FILE_NAME_CASE_STUDIES_GSLIDES = "National Campaigns_case studies" # From your image
FILE_NAME_CASE_STUDIES_PDF = "National_Campaigns_case_studies.pdf" # New or alternative
FILE_NAME_PHYSICAL_CAMPAIGNS_GSHEET = "Physical_campaigns_live_sheet" # From your image
FILE_NAME_DIGITAL_CAMPAIGNS_GSHEET = "Digital_Campaigns_live_sheet" # From your image (note "live sheet")
FILE_NAME_COM_DATA_GSHEET = "NoBroker_Overall_Data" # From your image
FILE_NAME_NBH_PREVIOUS_MEETINGS_GSHEET = "NBH_previous_meetings_updated"
FILE_NAME_LATEST_CASE_STUDIES_GSHEET = "Consolidated Case Studies - Master"

def parse_names_from_cell_helper(cell_value_str):
    """
    Intelligently parses a string from a spreadsheet cell to extract a set of cleaned, lowercase names.
    Handles two primary formats:
    1. A string representation of a list (e.g., "['name1@example.com', 'name2']").
    2. A simple delimited string (e.g., "Name One, Name Two & Name Three").
    """
    names = set()
    if not cell_value_str or str(cell_value_str).strip().lower() == 'n/a':
        return names

    # Ensure we are working with a string and clean it up
    cell_value_str = str(cell_value_str).strip()
    potential_names = []

    # --- NEW LOGIC: Check if the string is formatted like a list ---
    if cell_value_str.startswith('[') and cell_value_str.endswith(']'):
        # Use regex to find all content within single or double quotes
        # This is safer than using eval()
        extracted_items = re.findall(r"[\'\"](.*?)[\'\"]", cell_value_str)
        for item in extracted_items:
            # If the item is an email, just take the part before the @
            if '@' in item:
                potential_names.append(item.split('@')[0])
            else:
                potential_names.append(item)
    else:
        # --- FALLBACK LOGIC: Handle simple, delimited strings ---
        # Remove content in parentheses (e.g., "(NoBrokerHood)")
        cleaned_cell = re.sub(r'\s*\([^)]*\)', '', cell_value_str)
        cleaned_cell = cleaned_cell.replace('*', '').strip() # Remove asterisks

        # Split by common delimiters
        potential_names = re.split(r'\s*[,;/&\n]\s*|\s+\band\b\s+|\s+\bwith\b\s+', cleaned_cell)

    # --- Common cleaning process for all extracted name parts ---
    for name_part in potential_names:
        final_name = name_part.strip().lower()
        # Filter for meaningful names and exclude common role descriptors
        if final_name and len(final_name) > 2 and \
           "nbh sales" not in final_name and \
           "brand representative" not in final_name and \
           "nobrokerhood" not in final_name and \
           "stay vista" not in final_name:
            names.add(final_name)
            
    return names

# --- ADD THESE NEW HELPER FUNCTIONS NEAR THE TOP OF YOUR SCRIPT ---

def normalize_attendee_name(name_str):
    """
    Takes a raw name string and converts it into a standardized set of name parts.
    - "Shubham Chandrakant Dakhane" -> {'shubham', 'chandrakant', 'dakhane'}
    - "shubham.chandrakant" -> {'shubham', 'chandrakant'}
    - "trisha.bagchi7" -> {'trisha', 'bagchi'}
    - "mary" -> {'mary'}
    """
    if not isinstance(name_str, str) or not name_str.strip():
        return set()
    
    # Lowercase and replace common delimiters with spaces
    processed_name = name_str.lower().replace('.', ' ').replace('_', ' ')
    
    # Remove all characters that are not letters or spaces
    processed_name = re.sub(r'[^a-z\s]', '', processed_name)
    
    # Split into parts and filter out any empty strings resulting from multiple spaces
    name_parts = {part for part in processed_name.split() if part}
    
    return name_parts

def find_common_attendees(attendee_set_1_raw, attendee_set_2_raw):
    """
    Compares two sets of raw name strings and finds common individuals
    using a flexible, normalization-based approach.
    Returns a list of the matched raw names from the first set.
    """
    # Normalize all names in both sets
    # Each item in these lists will be a set of name parts, e.g., [{'shubham', 'chandrakant'}, {'trisha', 'bagchi'}]
    normalized_attendees_1 = [normalize_attendee_name(name) for name in attendee_set_1_raw]
    normalized_attendees_2 = [normalize_attendee_name(name) for name in attendee_set_2_raw]

    common_attendees_raw_names = []
    
    # Keep track of which attendees from set 2 have already been matched to avoid double counting
    matched_indices_in_set_2 = set()

    for i, norm_set_1 in enumerate(normalized_attendees_1):
        if not norm_set_1:
            continue
        
        for j, norm_set_2 in enumerate(normalized_attendees_2):
            if j in matched_indices_in_set_2 or not norm_set_2:
                continue

            # --- NEW, MORE ROBUST CORE LOGIC ---
            # A match occurs if:
            # 1. The name sets are identical (e.g., {'john', 'doe'} == {'john', 'doe'}).
            # 2. One set is a complete subset of the other (e.g., {'john'} is a subset of {'john', 'doe'}).
            # 3. There is a non-trivial intersection (e.g., {'john', 'd'} and {'john', 'doe'} intersect on 'john').
            #    We add the len > 1 check to avoid matching on single initials like 'a'.

            intersection = norm_set_1.intersection(norm_set_2)
            is_match = False
            
            if norm_set_1 == norm_set_2:
                is_match = True
            elif norm_set_1.issubset(norm_set_2) or norm_set_2.issubset(norm_set_1):
                is_match = True
            elif intersection and any(len(name_part) > 1 for name_part in intersection):
                is_match = True

            if is_match:
                common_attendees_raw_names.append(list(attendee_set_1_raw)[i])
                matched_indices_in_set_2.add(j)
                break # Match found, move to the next person in set 1

    return common_attendees_raw_names

# --- END OF NEW HELPER FUNCTIONS ---
# ========== NEW FUNCTION: Smart & Strict Brand Matching (For Escalations) ==========
def is_brand_match(brand1, brand2):
    """
    Strict but smart brand matching for Escalation Emails.
    Matches 'Amazon' == 'Amazon India', 'Swiggy' == 'Swiggy Instamart'.
    Prevents 'Times' == 'Times OOH', 'ia' == 'Epigamia', or 'zee' == 'sbzee'.
    """
    if not brand1 or not brand2:
        return False
        
    b1 = str(brand1).lower().strip()
    b2 = str(brand2).lower().strip()
    
    if not b1 or not b2 or b1 == 'unknown' or b2 == 'unknown':
        return False
        
    # 1. Exact direct match (Fastest check)
    if b1 == b2:
        return True
        
    def clean_brand(b):
        # Replace special characters with space
        b = re.sub(r'[^a-z0-9]', ' ', b)
        
        # Fluff words that don't change the core brand identity
        words_to_ignore = {
            'india', 'pvt', 'ltd', 'private', 'limited', 'inc', 'corp', 
            'corporation', 'llc', 'the', 'group', 'co', 'company', 'brand', 'brands',
            'instamart', 'pay', 'fresh', 'digital', 'global', 'ventures', 'enterprise'
        }
        
        # Split into words, filter out ignored words
        words = b.split()
        cleaned =[w for w in words if w not in words_to_ignore]
        
        return ' '.join(cleaned).strip()

    clean_b1 = clean_brand(b1)
    clean_b2 = clean_brand(b2)
    
    # 2. Strict Exact Match on the cleaned core brand name
    if clean_b1 and clean_b2 and clean_b1 == clean_b2:
        return True
        
    return False
# ========== UPGRADED FUNCTION 1: Search for LinkedIn & Recent Posts (LIVE SERPER MODE) ==========
def search_attendee_intel(email, raw_name, company_name, gemini_llm_client):
    """
    Profiles attendees using Serper.dev across all production meetings.
    Handles personal profiles vs activity posts dynamically.
    """
    if not gemini_llm_client:
        return None
    
    email_prefix = email.split('@')[0] if '@' in email else email
    domain = email.split('@')[1] if '@' in email else ""

    search_query = f"{raw_name or email_prefix} {company_name or domain} India LinkedIn"
    print(f"    🔍 [Attendee Search] Querying Serper: '{search_query}'")
    search_context = execute_serper_search_api(search_query, num_results=4)
    
    prompt = f"""
    You are an expert OSINT researcher profiling a meeting attendee.
    Email: '{email}' | Company Name: '{company_name}' in India.
    Raw Name from calendar: '{raw_name}'

    Here are real-time search results:
    ---
    {search_context}
    ---

    Task 1: Deduce their correct full name.
    Task 2: Extract their personal profile URL (must contain 'linkedin.com/in/').
    Task 3: Extract any recent post/activity link from the search results (usually contains 'linkedin.com/posts/').

    RULES:
    1. Only return URLs present in the search data. Do not make up links.
    2. Do not return company pages, only personal profiles or their direct post URLs.
    3. If no matching LinkedIn profile is found, return null.

    Return ONLY a JSON object:
    {{
        "inferred_name": "Full Name",
        "linkedin_url": "https://www.linkedin.com/in/... or null",
        "recent_post_url": "https://www.linkedin.com/posts/... or null",
        "post_context": "Short description of post topic or null"
    }}
    """
    config = types.GenerateContentConfig(temperature=0.0, response_mime_type="application/json")
    try:
        response = gemini_llm_client.models.generate_content(model="gemini-2.5-flash", contents=prompt, config=config)
        data = json.loads(response.text.strip())
        return data
    except Exception as e:
        print(f"  Error parsing Serper response for {email}: {e}")
        return {"inferred_name": raw_name.title() or email_prefix.title(), "linkedin_url": None, "recent_post_url": None, "post_context": None}


# ========== UPGRADED FUNCTION 2: Get LinkedIn & Posts for All Attendees ==========
def get_brand_attendees_linkedin_info(brand_attendees_list, brand_name, gemini_llm_client):
    """
    For each brand attendee, search for their LinkedIn profile and recent activity.
    Returns a list with clean names, LinkedIn URLs, and Post URLs added.
    """
    attendees_with_intel = []
    
    for attendee in brand_attendees_list:
        attendee_name = attendee.get('name', '')
        attendee_email = attendee.get('email', '')
        
        print(f"    🔍 Searching OSINT Intel for: {attendee_email} at {brand_name}")
        
        # Call optimized Serper search function
        intel_data = search_attendee_intel(attendee_email, attendee_name, brand_name, gemini_llm_client)
        
        if not intel_data:
            intel_data = {"inferred_name": attendee_name, "linkedin_url": None, "recent_post_url": None, "post_context": None}

        attendees_with_intel.append({
            'name': intel_data.get('inferred_name', attendee_name),
            'email': attendee_email,
            'linkedin_url': intel_data.get('linkedin_url') or '(LinkedIn Not Verified)',
            'recent_post_url': intel_data.get('recent_post_url'),
            'post_context': intel_data.get('post_context')
        })
        
        # Reduced buffer to speed up real-time execution while respecting rate limits
        print("    ⏳ Waiting 1.5s between attendee queries...")
        time.sleep(1.5)
    
    return attendees_with_intel

# ========== NEW FUNCTION 3: Find Potential Key Contacts (LIVE SERPER MODE) ==========
def find_potential_key_contacts(brand_name, gemini_llm_client):
    """
    Finds 2-3 current execution-level brand leaders at the company in India using Serper.
    """
    if not gemini_llm_client:
        return []

    # Strict search operators prioritizing mid-level roles directly at the target brand
    search_query = f'site:linkedin.com/in/ "{brand_name}" ("Brand Manager" OR "Associate Brand Manager" OR "Marketing Manager" OR "Campaign Manager")'
    print(f"    🔍 [Key Contacts Search] Querying Serper: {search_query}")
    search_context = execute_serper_search_api(search_query, num_results=5)

    discovery_prompt = f"""
    You are an expert executive search strategist mapping key decision-makers.
    Brand: {brand_name} (India)

    Review this organic web search data:
    ---
    {search_context}
    ---

    Task: Identify 2-3 marketing or brand leaders currently in India.
    RULES:
    1. Extract their full name, exact title, and verified personal profile link (must contain 'linkedin.com/in/').
    
    Return ONLY this JSON format:
    {{
      "contacts": [
        {{"name": "Full Name", "title": "Job Title", "reasoning": "Coordinates partnerships", "linkedin_url": "Profile URL or null"}}
      ]
    }}
    """
    config = types.GenerateContentConfig(temperature=0.1, response_mime_type="application/json")
    try:
        response = gemini_llm_client.models.generate_content(model="gemini-2.5-flash", contents=discovery_prompt, config=config)
        contacts_data = json.loads(response.text.strip())
        return contacts_data.get("contacts", [])
    except Exception as e:
        print(f"    Error parsing key contacts from Serper: {e}")
        return []

class Industry(enum.Enum):
    FMCG = "FMCG"
    AUTOMOTIVE_AND_TRANSPORT = "Automotive & Transportation"
    MEMBERSHIP_AND_LOCAL_SERVICES="Membership & Local Services"
    MARKETING_ADVERTISING_AND_MEDIA="Marketing, Advertising & Media"
    APPAREL_AND_FASHION="Apparel & Fashion"
    FOOD_AND_BEVERAGE="Food & Beverage"
    HEALTHCARE="Healthcare"
    FINANCE_AND_FINTECH="Finance & Fintech"
    BEAUTY_AND_PERSONAL_CARE="Beauty & Personal Care"
    JEWELLERY="Jewellery"
    REAL_ESTATE_AND_CONSTRUCTION="Real Estate & Construction"
    ENERGY_RENEWABLES_AND_MINING="Energy, Renewables & Mining"
    WELLNESS_AND_FITNESS="Wellness & Fitness"
    EDUCATION_AND_TRAINING="Education & Training"
    HOME_GOODS_AND_ELECTRONINCS="Home Goods & Electronics"
    HOSPITALITY_AND_TRAVEL="Hospitality & Travel"
    TECHNOLOGY_AND_BUSINESS_SERVICES="Technology & Business Services"
    E_COMMERCE="E-Commerce"
    RETAIL="Retail"
    PETS_AND_PETS_SERVICES="Pets & Pet Services"
    GAMING="Gaming"
    LOGISTICS_AND_WAREHOUSING="Logistics & Warehousing"
    OTHER_UNKNOWN="Other / Unknown"
    MANUFACTURING_AND_INDUSTRIAL="Manufacturing & Industrial"
    QUICk_COMMERCE="Quick Commerce"
    PHARMA = "Pharma"
    OTT = "OTT"



# Master list of Media & Advertising Agencies
KNOWN_AGENCIES_LIST = [
    "wpp", "wpp media", "groupm", "mindshare", "wavemaker", "essencemediacom", 
    "publicis", "publicis media", "havas", "havas media", "pivotroots", 
    "omd", "omc", "omnicom", "omnicom media group", "starcom", "zenith", 
    "dentsu", "lyxel & flamingo", "lyxelandflamingo", "ls digital", "lsdigital", 
    "madison", "madison media", "madison ooh", "interactive avenues", 
    "hiveminds", "oap india", "konnect services"
]

class Brand_Details(BaseModel):
    brand_name: str
    industry: Industry
    sub_category_keywords: list[str]
    is_agency_meeting: bool = False
    agency_name: str = "N/A"
    
    class Config:
        use_enum_values = True  # Use enum values instead of names in JSON output

Allowed_Industries = [industry.value for industry in Industry]

BRAND_EXTRACTION_PROMPT_TEMPLATE = """
You are an expert administrative assistant working for NoBrokerHood (NBH) responsible for parsing meeting titles between NBH and external companies.

Target Title: "{MEETING_TITLE}"

Your task is to analyze the meeting title and extract information about the target company or brand.

CRITICAL RULES FOR AGENCY MEETINGS vs DIRECT BRAND MEETINGS:
1. **Identify Media/Advertising Agencies**:
   Known Agencies in India include: WPP, GroupM, Mindshare, Wavemaker, EssenceMediacom, Publicis, Havas, PivotRoots, OMD, Omnicom, Starcom, Zenith, Dentsu, Lyxel & Flamingo, LS Digital, Madison, Interactive Avenues, Hiveminds, OAP India.

2. **Rule A - Agency Title WITH Client Brand** (e.g., "WPP (Prime Video) X NBH", "Dentsu - Maruti Suzuki x NBH"):
   - `is_agency_meeting`: true
   - `agency_name`: The Agency Name (e.g., "WPP" or "Dentsu")
   - `brand_name`: The specific client brand inside parentheses or title (e.g., "Prime Video" or "Maruti Suzuki")
   - `industry`: Infer industry of the client brand (e.g., "OTT" or "Automotive & Transportation")

3. **Rule B - Agency Title WITHOUT Any Specific Client Brand** (e.g., "In-Person | Dentsu x NBH", "WPP || NoBrokerHood", "Publicis x NBH"):
   - `is_agency_meeting`: true
   - `agency_name`: The Agency Name (e.g., "Dentsu", "WPP", "Publicis")
   - `brand_name`: DO NOT set brand_name to the Agency itself! Instead, identify a major, highly-recognized consumer client brand in India handled by this agency's portfolio (e.g., For Dentsu: "Toyota" or "Maruti Suzuki"; For WPP/GroupM/Mindshare: "Cadbury" or "Tata Motors"; For Publicis: "L'Oreal" or "Nestle"; For Havas: "Reckitt" or "Citroen").
   - `industry`: Infer industry of that selected portfolio client brand.

4. **Rule C - Direct Brand Meeting** (e.g., "Harpic x NBH", "GIVA Digital discussion"):
   - `is_agency_meeting`: false
   - `agency_name`: "N/A"
   - `brand_name`: The primary brand being met (e.g., "Harpic", "Giva").

5. **sub_category_keywords**: Provide 4 to 6 specific keywords/competitor brands in India for the `brand_name`.
6. **industry**: Choose strictly from **Allowed_Industries**: {Allowed_Industries}

Return ONLY a JSON object matching this schema:
{{
  "brand_name": "...",
  "industry": "...",
  "sub_category_keywords": ["..."],
  "is_agency_meeting": true/false,
  "agency_name": "..."
}}
"""

def get_brand_details_from_title_with_llm(gemini_llm_client, meeting_title):
    """
    Uses LLM to extract brand details, dynamically detecting agency meetings 
    and selecting a representative portfolio brand when only an agency is named.
    """
    default_response = {
        "brand_name": "Unknown Brand",
        "industry": "Unknown",
        "sub_category_keywords": [],
        "is_agency_meeting": False,
        "agency_name": "N/A"
    }
    if not gemini_llm_client:
        print("  LLM model not available for brand extraction.")
        return default_response

    prompt = BRAND_EXTRACTION_PROMPT_TEMPLATE.format(
        MEETING_TITLE=meeting_title, 
        Allowed_Industries=Allowed_Industries
    )

    grounding_tool = types.Tool(google_search=types.GoogleSearch())
    config = types.GenerateContentConfig(tools=[grounding_tool])

    try:
        raw_text = ""
        response = gemini_llm_client.models.generate_content(
            model="gemini-2.5-flash", 
            contents=prompt, 
            config=config
        )
        raw_text = response.candidates[0].content.parts[0].text

        cleaned_json_str = re.sub(r'```json\s*|\s*```', '', raw_text).strip()
        data = json.loads(cleaned_json_str)

        if "brand_name" in data and "industry" in data:
            extracted_brand = str(data.get("brand_name", "")).strip().lower()
            
            # Python Safety Net: If LLM still returned an agency name as brand_name
            is_agency_hit = any(agency in extracted_brand for agency in KNOWN_AGENCIES_LIST)
            
            if is_agency_hit:
                print(f"  ⚠️ LLM extracted agency '{extracted_brand}' as brand. Overriding with portfolio brand logic...")
                data["is_agency_meeting"] = True
                data["agency_name"] = extracted_brand.title()
                
                # Direct portfolio mapping for fallback
                agency_portfolio_fallback = {
                    "dentsu": ("Maruti Suzuki", "Automotive & Transportation", ["cars", "suv", "hyundai", "tata motors"]),
                    "wpp": ("Cadbury", "FMCG", ["chocolates", "confectionery", "nestle", "amul"]),
                    "wpp media": ("Cadbury", "FMCG", ["chocolates", "confectionery", "nestle", "amul"]),
                    "groupm": ("Tata Motors", "Automotive & Transportation", ["cars", "ev", "mahindra", "hyundai"]),
                    "mindshare": ("Unilever", "FMCG", ["personal care", "home care", "p&g", "itc"]),
                    "wavemaker": ("L'Oreal", "Beauty & Personal Care", ["shampoo", "skincare", "maybelline", "lakme"]),
                    "essencemediacom": ("Google", "Technology & Business Services", ["search", "pixel", "tech", "android"]),
                    "publicis": ("Nestle", "Food & Beverage", ["coffee", "noodles", "maggie", "britannia"]),
                    "publicis media": ("Nestle", "Food & Beverage", ["coffee", "noodles", "maggie", "britannia"]),
                    "havas": ("Reckitt", "FMCG", ["dettol", "harpic", "lizol", "savlon"]),
                    "havas media": ("Reckitt", "FMCG", ["dettol", "harpic", "lizol", "savlon"]),
                    "madison": ("Godrej", "FMCG", ["soap", "home care", "dabur", "marico"]),
                    "madison media": ("Godrej", "FMCG", ["soap", "home care", "dabur", "marico"]),
                    "omd": ("Apple", "Home Goods & Electronics", ["iphone", "macbook", "samsung", "oneplus"]),
                    "omnicom": ("Apple", "Home Goods & Electronics", ["iphone", "macbook", "samsung", "oneplus"]),
                    "omnicom media group": ("Apple", "Home Goods & Electronics", ["iphone", "macbook", "samsung", "oneplus"]),
                    "pivotroots": ("Amazon", "E-Commerce", ["shopping", "delivery", "prime", "flipkart"]),
                    "starcom": ("Samsung", "Home Goods & Electronics", ["smartphones", "tv", "appliances", "lg"]),
                    "zenith": ("Disney+ Hotstar", "OTT", ["streaming", "movies", "ipl", "netflix"])
                }
                
                matched = None
                for key, val in agency_portfolio_fallback.items():
                    if key in extracted_brand:
                        matched = val
                        break
                
                if not matched:
                    # Generic high-performing fallback for any unspecified agency
                    matched = ("Cadbury", "FMCG", ["chocolates", "confectionery", "nestle", "amul"])

                data["brand_name"] = matched[0]
                data["industry"] = matched[1]
                data["sub_category_keywords"] = matched[2]

            if not data["brand_name"] or data["brand_name"].lower() == 'unknown':
                print(f"  LLM identified title '{meeting_title}' as ambiguous.")
                return default_response
            return data

        return default_response

    except Exception as e:
        print(f"⚠️ First pass failed ({e}), retrying with strict JSON schema…")
        try:
            cleanup_prompt = (
                "Reformat the text below into valid JSON matching this schema:\n\n"
                f"{json.dumps(Brand_Details.model_json_schema(), indent=2)}\n\n"
                f"{raw_text}"
            )
            cleanup_config = types.GenerateContentConfig(
                response_mime_type="application/json",
                response_schema=Brand_Details
            )
            retry = gemini_llm_client.models.generate_content(
                model="gemini-2.5-flash",
                contents=cleanup_prompt,
                config=cleanup_config
            )
            parsed: Brand_Details = retry.parsed
            return parsed.model_dump()
        except Exception as e2:
            print(f"❌ Retry failed: {e2}")
            return default_response


# --- Google Authentication and Service Building ---
def get_google_service(service_name, version, scopes_list, token_filename_base_for_local_storage): # Changed last param name for clarity
    creds = None
    # Construct the specific local token filename (e.g., token_brandvmeet_calendar.json)
    # This is used for local development fallback and saving tokens locally.
    local_token_file_path = token_filename_base_for_local_storage

    # --- Attempt 1: Load from specific environment variable for the service (for CI) ---
    # Construct the expected environment variable name, e.g., GOOGLE_TOKEN_JSON_CALENDAR
    token_env_var_name = f"GOOGLE_TOKEN_JSON_{service_name.upper()}"
    token_json_string_from_env = os.getenv(token_env_var_name)

    if os.getenv('CI') == 'true' and token_json_string_from_env: # Check 'CI' env var and if token string exists
        print(f"CI environment detected. Attempting to load credentials for {service_name} from env var: {token_env_var_name}")
        try:
            token_info = json.loads(token_json_string_from_env)
            # The token_info from your stored JSON should contain client_id, client_secret, and refresh_token,
            # which are needed for the refresh mechanism by the Credentials object.
            creds = Credentials.from_authorized_user_info(token_info, scopes_list)
            print(f"Successfully loaded credentials for {service_name} from environment variable.")
        except json.JSONDecodeError as e:
            print(f"Error decoding JSON from {token_env_var_name} for {service_name}: {e}")
            creds = None
        except Exception as e:
            print(f"Generic error loading credentials for {service_name} from environment variable {token_env_var_name}: {e}")
            creds = None # Fallback

    # --- Attempt 2: Load from local token file (for local development or if CI load failed) ---
    if not creds and os.path.exists(local_token_file_path):
        print(f"Loading credentials for {service_name} from local file: {local_token_file_path}")
        try:
            creds = Credentials.from_authorized_user_file(local_token_file_path, scopes_list)
        except Exception as e:
            print(f"Error loading credentials from {local_token_file_path} for {service_name}: {e}")
            creds = None

    # --- Attempt 3 & 4: Refresh or Run Interactive Flow ---
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            print(f"Token for {service_name} is expired. Refreshing...")
            try:
                creds.refresh(Request())
                print(f"Token for {service_name} refreshed successfully.")
                # If refreshed successfully in a non-CI environment, save it back to the local file
                if not os.getenv('CI') == 'true':
                    try:
                        with open(local_token_file_path, 'w') as token_file:
                            token_file.write(creds.to_json())
                        print(f"Saved refreshed token for {service_name} to {local_token_file_path}")
                    except Exception as e_save:
                        print(f"Error saving refreshed token for {service_name} to {local_token_file_path}: {e_save}")
            except Exception as e_refresh:
                print(f"Error refreshing token for {service_name}: {e_refresh}")
                # If refresh fails in CI, it's a critical issue as there's no interactive fallback.
                if os.getenv('CI') == 'true':
                    print(f"FATAL: Token refresh failed for {service_name} in CI. Stored token might be invalid (e.g., revoked) or scopes changed.")
                    return None
                creds = None # Force re-authentication locally if refresh failed and not in CI
        
        # This block should ideally NOT be reached in a CI environment if the token env var is set correctly.
        if not creds:
            if os.getenv('CI') == 'true':
                print(f"FATAL: No valid credentials for {service_name} in CI and interactive flow is disabled.")
                print(f"       Ensure {token_env_var_name} secret is set correctly and contains valid token JSON.")
                return None # Critical failure in CI

            print(f"No valid credentials for {service_name}. Attempting interactive local server flow...")
            # Ensure CREDENTIALS_FILE ('credentials.json' with OAuth client_id/secret) exists for the flow
            if not os.path.exists(CREDENTIALS_FILE):
                print(f"FATAL: {CREDENTIALS_FILE} not found. Cannot run interactive auth flow for {service_name}.")
                return None
            try:
                flow = InstalledAppFlow.from_client_secrets_file(CREDENTIALS_FILE, scopes_list)
                creds = flow.run_local_server(port=0) # This opens a browser for user consent
                # Save the newly obtained credentials to the local token file for future runs
                with open(local_token_file_path, 'w') as token_file:
                    token_file.write(creds.to_json())
                print(f"Saved new token for {service_name} (from interactive flow) to {local_token_file_path}")
            except Exception as e_flow:
                print(f"Error during interactive auth flow for {service_name}: {e_flow}")
                return None
    
    if not creds: # Should not happen if all paths above are handled
        print(f"Ultimately failed to obtain credentials for {service_name}.")
        return None

    # --- Build the Google API Service ---
    try:
        service = build(service_name, version, credentials=creds)
        print(f"{service_name.capitalize()} service initialized successfully.")
        return service
    except HttpError as error:
        print(f'An HTTP error occurred building {service_name} service: {error}')
        return None
    except Exception as e: # Catch other potential errors during build
        print(f'A general error occurred building {service_name} service: {e}')
        return None
    
# --- Google Drive Functions ---
def list_files_in_gdrive_folder(drive_service, folder_id):
    # ... (Implementation from previous thought block - list files) ...
    # Ensure it handles empty folder_id
    if not folder_id or folder_id == "YOUR_GDRIVE_FOLDER_ID_HERE":
        print("Google Drive Folder ID for NBH data is not configured.")
        return []
    try:
        results = drive_service.files().list(
            q=f"'{folder_id}' in parents and trashed=false",
            # Include mimeType to decide how to read
            fields="nextPageToken, files(id, name, mimeType)"
        ).execute()
        items = results.get('files', [])
        return items
    except HttpError as error:
        print(f"An error occurred listing GDrive files for folder {folder_id}: {error}")
        return []

    

# THIS IS THE PRIMARY FUNCTION TO GET FILE DATA
def get_structured_gdrive_file_data(drive_service, sheets_service, file_id, file_name, mime_type):
    """
    Extracts and parses structured content from a Google Drive file based on its MIME type.
    
    Supports Google Slides (as PPTX), Microsoft PowerPoint, Google Sheets, Microsoft Excel, Google Docs, PDFs, and plain text files. Returns structured data as a list of dictionaries for presentations and spreadsheets, or as a string for text-based files. If parsing fails or the file type is unsupported, returns a descriptive error message.
    """

    print(f"    Attempting to read structured data for: {file_name} (MIME: {mime_type})")

    MIMETYPE_GOOGLE_SHEET = 'application/vnd.google-apps.spreadsheet'
    MIMETYPE_GOOGLE_DOC = 'application/vnd.google-apps.document'
    MIMETYPE_GOOGLE_PRESENTATION = 'application/vnd.google-apps.presentation'
    MIMETYPE_MS_EXCEL = 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    MIMETYPE_MS_POWERPOINT = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    MIMETYPE_PDF = 'application/pdf'

    try:
        # --- Google Slides (Presentations) ---
        if mime_type == MIMETYPE_GOOGLE_PRESENTATION:
            print(f"    Exporting Google Slides '{file_name}' as PPTX for slide-level parsing...")
            # Export Google Slides as PPTX
            request_pptx = drive_service.files().export_media(
                fileId=file_id,
                mimeType=MIMETYPE_MS_POWERPOINT # Export as PPTX
            )
            fh_pptx = io.BytesIO()
            downloader_pptx = MediaIoBaseDownload(fh_pptx, request_pptx)
            done_pptx = False
            while not done_pptx:
                status, done_pptx = downloader_pptx.next_chunk()
            fh_pptx.seek(0)

            try:
                prs = Presentation(fh_pptx)
                slides_data = []
                for i, slide in enumerate(prs.slides):
                    slide_text_runs = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame") and shape.text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.text:
                                        slide_text_runs.append(run.text)
                    slides_data.append({
                        "slide_number": i + 1,
                        "text": "\n".join(slide_text_runs),
                        "source_type": "Google Presentation (exported as PPTX)"
                    })
                if not slides_data and prs.slides: # Had slides but no text extracted
                     return f"Google Slides '{file_name}': Exported as PPTX, but no text content found in slides."
                elif not prs.slides: # Exported PPTX was empty
                    return f"Google Slides '{file_name}': Exported as PPTX, but it contained no slides."
                return slides_data                          
            except Exception as e_pptx_parse:
                print(f"    Warning: Failed to parse Google Slides '{file_name}' exported as PPTX: {e_pptx_parse}. Falling back to plain text export.")
                # Fallback: Export as plain text if PPTX export/parse fails
                request_text = drive_service.files().export_media(fileId=file_id, mimeType='text/plain')
                fh_text = io.BytesIO()
                downloader_text = MediaIoBaseDownload(fh_text, request_text)
                done_text = False
                while not done_text:
                    status, done_text = downloader_text.next_chunk()
                fh_text.seek(0)
                full_text = fh_text.read().decode('utf-8', errors='replace')
                if not full_text.strip():
                    return f"Google Slides '{file_name}': Fallback to plain text export, but no content found."
                return [{"slide_number": 1, "text": full_text, "source_type": "Google Presentation (all text fallback)"}]


        # --- Microsoft PowerPoint (.pptx) ---
        elif mime_type == MIMETYPE_MS_POWERPOINT: # Handles native PPTX files
            request = drive_service.files().get_media(fileId=file_id)
            fh = io.BytesIO()
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while not done: status, done = downloader.next_chunk()
            fh.seek(0)
            try:
                prs = Presentation(fh)
                slides_data = []
                for i, slide in enumerate(prs.slides):
                    slide_text_runs = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text_frame") and shape.text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.text: slide_text_runs.append(run.text)
                    slides_data.append({"slide_number": i + 1, "text": "\n".join(slide_text_runs), "source_type": "PPTX (native)"}) # Clarified source
                if not slides_data and prs.slides:
                    return f"PPTX File '{file_name}': Contained slides, but no text content found in them."
                elif not prs.slides:
                    return f"PPTX File '{file_name}': Contained no slides."
                return slides_data
            except Exception as e:
                return f"Could not parse native PPTX content for '{file_name}': {e}"
            

        # --- Google Sheets & Microsoft Excel ---
        elif mime_type == MIMETYPE_GOOGLE_SHEET or mime_type == MIMETYPE_MS_EXCEL:
            all_rows_data = []
            try:
                if mime_type == MIMETYPE_GOOGLE_SHEET:
                    spreadsheet = sheets_service.spreadsheets().get(spreadsheetId=file_id).execute()
                    for sheet_meta in spreadsheet.get('sheets', []):
                        sheet_title = sheet_meta['properties']['title']
                        # Read a significant number of rows, e.g., 1000, and all columns up to ZZ
                        range_str = f"'{sheet_title}'!A1:AZ5000" 
                        result = sheets_service.spreadsheets().values().get(
                            spreadsheetId=file_id, range=range_str
                        ).execute()
                        rows = result.get('values', [])
                        if rows:
                            header = rows[0]
                            for row_idx, row_values in enumerate(rows): # (row_idx is 0-based here)
                                all_rows_data.append({"sheet_name": sheet_title, 
                                                      "row_index": row_idx + 1, # 1-based for display
                                                      "header": header, 
                                                      "values": row_values,
                                                      "source_type": "Google Sheet"})
                elif mime_type == MIMETYPE_MS_EXCEL:
                    request = drive_service.files().get_media(fileId=file_id)
                    fh = io.BytesIO()
                    downloader = MediaIoBaseDownload(fh, request)
                    done = False
                    while not done: status, done = downloader.next_chunk()
                    fh.seek(0)
                    workbook = openpyxl.load_workbook(fh)
                    for sheet_title in workbook.sheetnames:
                        sheet = workbook[sheet_title]
                        header = [cell.value for cell in sheet[1]] # Assuming header is 1st row
                        for row_idx, row_obj in enumerate(sheet.iter_rows(min_row=1, max_col=50, max_row=1000, values_only=True)): # max_col to limit width
                            # row_idx is 0-based here from iter_rows(min_row=1)
                            all_rows_data.append({"sheet_name": sheet_title,
                                                  "row_index": row_idx + 1, # 1-based for display
                                                  "header": header, 
                                                  "values": list(row_obj), # Ensure it's a list
                                                  "source_type": "Excel Sheet"})
                if not all_rows_data: # If loops completed but no data (e.g. all sheets were empty)
                    return f"Spreadsheet file '{file_name}' ({mime_type}) processed, but no data rows found in any sheet."
                return all_rows_data
            except Exception as e:
                return f"Could not parse Spreadsheet content for '{file_name}' ({mime_type}): {e}"
        
        # --- Google Documents ---
        elif mime_type == MIMETYPE_GOOGLE_DOC:
            request = drive_service.files().export_media(fileId=file_id, mimeType='text/plain')
            fh = io.BytesIO()
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while not done: status, done = downloader.next_chunk()
            fh.seek(0)
            return fh.read().decode('utf-8', errors='replace') # Returns string

        # --- PDF Files ---
        elif mime_type == MIMETYPE_PDF:
            try:
                request = drive_service.files().get_media(fileId=file_id)
                fh = io.BytesIO()
                downloader = MediaIoBaseDownload(fh, request)
                done = False
                while not done:
                    status, done = downloader.next_chunk()
                fh.seek(0)
                pdf_text = ""
                with fitz.open(stream=fh.read(), filetype="pdf") as doc:
                    for page in doc:
                        pdf_text += page.get_text() + "\n"
                if not pdf_text.strip():
                    return f"PDF File: '{file_name}'. No extractable text found."
                return pdf_text
            except Exception as e:
                return f"Could not parse PDF content: {e}" 

        # --- Plain Text Files ---
        elif mime_type.startswith('text/'):
            request = drive_service.files().get_media(fileId=file_id)
            fh = io.BytesIO()
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while not done: status, done = downloader.next_chunk()
            fh.seek(0)
            return fh.read().decode('utf-8', errors='replace') # Returns string
            
        else:
            return f"File type ({mime_type}) for '{file_name}' not configured for structured data extraction. Its name and existence might be relevant." # Returns string

    except HttpError as error:
        return f"An HTTP error occurred reading GDrive file {file_name} (ID: {file_id}): {error}" # Returns string
    except Exception as e:
        return f"A general error occurred reading GDrive file {file_name} (ID: {file_id}): {e}" # Returns string


def summarize_file_content_with_gemini(gemini_llm_client, file_name, mime_type, file_content):
    """
    Uses Gemini LLM to summarize the content of a file for inclusion in the meeting brief.
    """
    if not gemini_llm_client:
        return "Error: Gemini model not available for summarization."

    prompt = (
        f"Summarize the following content from the file '{file_name}' (type: {mime_type}) in 5-10 concise bullet points, "
        "focusing on key facts, data, or insights that would be useful for a marketing/sales meeting. "
        "Do not include generic statements. If the content is not relevant, say 'No relevant content found.'\n\n"
        f"---\n{file_content}\n---"
    )
    try:
        response = gemini_llm_client.models.generate_content(model="gemini-2.5-flash",contents=prompt)
        if response and response.candidates and response.candidates[0].content.parts:
            summary = response.candidates[0].content.parts[0].text.strip()
            return summary
        else:
            return "No summary generated."
    except Exception as e:
        print(f"Error during Gemini summarization: {e}")
        return f"Error: Exception during Gemini summarization: {e}"

# ==============================================================================
# STRICT EXACT MATCH HELPER: ONLY Exact Brand OR Exact Industry (No broad fallbacks)
# ==============================================================================
# ==============================================================================
# STRICT EXACT MATCH HELPER: Deep Scan for 2025/2026 Data
# ==============================================================================
def extract_strict_campaigns_and_case_studies(file_data_obj, fname, brand_clean, strict_keywords, sub_category_keywords=None, target_cities=None, target_depts=None, email_to_geo_map=None):
    """
    Scans sheet from BOTTOM (Newest) to TOP (Oldest). Filters ONLY for 2025 and 2026 executed data.
    Prioritizes campaigns matching BOTH City and Department, or just City for Case Studies.
    """
    if sub_category_keywords is None: sub_category_keywords = []
    if target_cities is None: target_cities = set()
    if target_depts is None: target_depts = set()
    if email_to_geo_map is None: email_to_geo_map = {}

    matches_brand_priority, matches_brand_other = [], []
    matches_sub_category_priority, matches_sub_category_other = [], []
    matches_strict_priority, matches_strict_other = [], []
    
    if not isinstance(file_data_obj, list) or not file_data_obj: 
        return []
    
    header_vals = file_data_obj[0].get("header",[]) if isinstance(file_data_obj[0], dict) else []
    brand_col, ind_col, date_col, city_col, email_col = -1, -1, -1, -1, -1

    if header_vals:
        lower_h = [str(h).strip().lower() for h in header_vals]
        for idx, h in enumerate(lower_h):
            if "brand" in h: brand_col = idx
            if any(x in h for x in ["industry", "category", "vertical"]): ind_col = idx
            if any(x in h for x in ["year", "date", "timestamp", "tentative"]): date_col = idx
            if "city" in h: city_col = idx
            if "email" in h: email_col = idx

    data_rows = file_data_obj[1:] if len(file_data_obj) > 1 else []
    if not data_rows: return []

    valid_sub_kws = [k.strip().lower() for k in sub_category_keywords if k and len(k.strip()) > 2]
    valid_industry_kws = [k.strip().lower() for k in strict_keywords if k and len(k.strip()) > 2]

    for row_info in reversed(data_rows):
        vals = row_info.get('values', [])
        if not vals: continue
        
        row_brand = str(vals[brand_col]).strip().lower() if brand_col != -1 and len(vals) > brand_col else ""
        row_ind = str(vals[ind_col]).strip().lower() if ind_col != -1 and len(vals) > ind_col else ""
        
        if not row_brand or row_brand in ['nan', 'none', '', 'n/a', 'unknown']:
            continue

        # --- GEO EXTRACTION & MATCHING LOGIC ---
        row_city_raw = ""
        row_dept_raw = ""
        
        # 1. Direct City Column (Case Studies)
        if city_col != -1 and len(vals) > city_col:
            row_city_raw = str(vals[city_col]).strip().lower()
            
        # 2. Email Mapping (Campaigns)
        if email_col != -1 and len(vals) > email_col:
            row_email = str(vals[email_col]).strip().lower()
            geo_info = email_to_geo_map.get(row_email, {})
            if geo_info:
                if not row_city_raw: row_city_raw = geo_info.get('city', '')
                row_dept_raw = geo_info.get('dept', '')
                
        # 3. Evaluate Matches
        is_city_match = target_cities and row_city_raw and any(tc in row_city_raw for tc in target_cities)
        is_dept_match = target_depts and row_dept_raw and any(td in row_dept_raw for td in target_depts)
        
        # 4. Generate Tag
        tag = ""
        is_priority = False
        if is_city_match and is_dept_match and row_dept_raw: # Campaign Match
            tag = " [📍 SAME CITY & DEPT MATCH]"
            is_priority = True
        elif is_city_match: # Case Study Match (or Campaign with just City)
            tag = " [📍 SAME CITY MATCH]"
            is_priority = True

        row_items = []
        for i in range(len(header_vals)):
            if i < len(vals):
                val = str(vals[i]).strip()
                if val and val.lower() not in ['nan', 'none', '', 'n/a']:
                    row_items.append(f"{header_vals[i]}: {val}")
        
        entry = " | ".join(row_items)
        if len(entry) < 15: continue
        
        if "2025" not in entry and "2026" not in entry and "2027" not in entry: continue
        
        entry_with_tag = entry + tag
        
        # Priority 1: Exact Brand Match
        if brand_clean and len(brand_clean) > 2 and row_brand and len(row_brand) > 2:
            if (brand_clean in row_brand) or (row_brand in brand_clean):
                if is_priority and len(matches_brand_priority) < 10:
                    matches_brand_priority.append(entry_with_tag)
                elif not is_priority and len(matches_brand_other) < 10: 
                    matches_brand_other.append(entry_with_tag)
                continue 
        
        # Priority 2: Sub-Category / Competitor Match
        entry_lower = entry.lower()
        if valid_sub_kws and any(k in entry_lower for k in valid_sub_kws):
            if is_priority and len(matches_sub_category_priority) < 30:
                matches_sub_category_priority.append(entry_with_tag)
            elif not is_priority and len(matches_sub_category_other) < 30:
                matches_sub_category_other.append(entry_with_tag)
            continue
            
        # Priority 3: Broad Industry Match 
        if valid_industry_kws:
            if (ind_col != -1 and row_ind and any(k in row_ind for k in valid_industry_kws)) or \
               (brand_col != -1 and row_brand and any(k in row_brand for k in valid_industry_kws)):
                if is_priority and len(matches_strict_priority) < 50:
                    matches_strict_priority.append(entry_with_tag)
                elif not is_priority and len(matches_strict_other) < 50: 
                    matches_strict_other.append(entry_with_tag)

    # Return Logic: Prioritize Tagged matches to the top of the list for the LLM
    final_output = []
    
    if matches_brand_priority or matches_brand_other:
        final_output.append("**Exact Brand Matches Found:**")
        final_output += matches_brand_priority + matches_brand_other
        
    if matches_sub_category_priority or matches_sub_category_other:
        final_output.append("**Highly Relevant (Competitor/Sub-Category) Campaigns:**")
        final_output += matches_sub_category_priority + matches_sub_category_other
        
    if matches_strict_priority or matches_strict_other:
        # Re-labeled to signal to the LLM that these are broad fallbacks and must not be matched to different product categories
        final_output.append("**Broad Industry Fallbacks (Low Relevance - Match ONLY if product category is identical):**")
        final_output += matches_strict_priority + matches_strict_other
    
    return final_output

# ==============================================================================
# MAIN FUNCTION: Powered by Strict Exact Mapping
# ==============================================================================
def get_internal_nbh_data_for_brand(drive_service, sheets_service, gemini_llm_client, 
                                    current_target_brand_name, target_brand_industry, current_meeting_data, 
                                    EXCLUDED_NBH_PSEUDO_NAMES_FOR_FOLLOWUP, AGENT_EMAIL, master_sheet_id, email_to_geo_map=None):
    
    if email_to_geo_map is None: email_to_geo_map = {}
    print(f"Fetching and processing internal NBH data for target brand '{current_target_brand_name}'...")
    
    # NEW: Determine Target Cities and Departments for Attendees
    target_cities = set()
    target_depts = set()
    for att in current_meeting_data.get('nbh_attendees', []):
        att_email = att.get('email', '').lower()
        if att_email in email_to_geo_map:
            geo_info = email_to_geo_map[att_email]
            if geo_info.get('city'): target_cities.add(geo_info['city'])
            if geo_info.get('dept'): target_depts.add(geo_info['dept'])
            
    print(f"    🎯 Target Geo for Attendees -> Cities: {target_cities}, Depts: {target_depts}")
    
    history_context_str = ""
    data_buckets = {"physical_campaigns": [], "digital_campaigns": [], "case_studies":[], "general_docs":[]}
    is_overall_direct_follow_up = False
    has_other_past_interactions = False 
    condensed_past_meetings_for_alert =[]

    # --- 1. MASTER SHEET LOGIC (History) ---
    current_meeting_date_obj = current_meeting_data.get('start_time_obj')
    if isinstance(current_meeting_date_obj, datetime.datetime):
        current_meeting_date_only = current_meeting_date_obj.date()
    else:
        current_meeting_date_only = datetime.date.today()

    current_nbh_tokens = set()
    for att in current_meeting_data.get('nbh_attendees',[]):
        if att.get('email'): current_nbh_tokens.add(att['email'].lower().split('@')[0].strip()) 
        if att.get('name'):
             parts = att['name'].lower().split()
             for p in parts: 
                 if len(p) > 2: current_nbh_tokens.add(p)

    try:
        header_req = sheets_service.spreadsheets().values().get(spreadsheetId=master_sheet_id, range="Meeting_data!A1:AZ1").execute()
        headers = header_req.get('values', [])[0]
        lower_headers =[str(h).strip().lower() for h in headers]
        
        try:
            col_brand = lower_headers.index("brand name")
            col_date = lower_headers.index("meeting date")
            col_discussion = lower_headers.index("key discussion points")
            col_actions = lower_headers.index("action items")
            col_nbh_attendees = lower_headers.index("nobroker attendees")
        except:
             history_context_str = "## PREVIOUS MEETING INTELLIGENCE: NONE (Fresh Meeting)\n"
        else:
            data_req = sheets_service.spreadsheets().values().get(spreadsheetId=master_sheet_id, range="Meeting_data!A2:AZ").execute()
            data_rows = data_req.get('values', [])
            found_meetings =[]
            target_clean = current_target_brand_name.lower().strip()
            
            current_meeting_id = current_meeting_data.get('id', '')

            for row in data_rows:
                if len(row) <= col_brand: continue
                
                row_meeting_id = str(row[0]).strip() if len(row) > 0 else ""
                if row_meeting_id and row_meeting_id == current_meeting_id:
                    continue

                sheet_brand = str(row[col_brand]).strip().lower()
                
                if not sheet_brand or sheet_brand == 'unknown' or not target_clean or target_clean == 'unknown':
                    continue
                
                if is_brand_match(target_clean, sheet_brand):
                    row_date_str = str(row[col_date]) if len(row) > col_date else ""
                    prev_nbh_raw = str(row[col_nbh_attendees]).lower() if len(row) > col_nbh_attendees else ""
                    is_attendee_match = any(token in prev_nbh_raw for token in current_nbh_tokens)
                    
                    meeting_info = {
                        "date": row_date_str,
                        "discussion": row[col_discussion] if len(row) > col_discussion else "N/A",
                        "actions": row[col_actions] if len(row) > col_actions else "None",
                        "nbh_team": prev_nbh_raw
                    }

                    if is_attendee_match: found_meetings.append(meeting_info)
                    else:
                        has_other_past_interactions = True
                        condensed_past_meetings_for_alert.append({"date": row_date_str, "discussion_summary": "Different Team - Content Hidden", "nbh_team": prev_nbh_raw})
            
            if found_meetings:
                is_overall_direct_follow_up = True
                top = found_meetings[0]
                history_context_str = f"## PREVIOUS MEETING INTELLIGENCE (MATCHED)\n**Last Meeting Date:** {top['date']}\n**Last Discussion:** {top['discussion']}\n**Last Actions:** {top['actions']}\n"
            else:
                 history_context_str = "## PREVIOUS MEETING INTELLIGENCE: NONE (Fresh Meeting)\n"

    except Exception as e:
        print(f"    Error reading Master Sheet: {e}")
        history_context_str = "## PREVIOUS MEETING INTELLIGENCE: Data access error.\n"

    # --- 2. EXACT INDUSTRY MAPPING (Strict Matching Only) ---
    STRICT_INDUSTRY_MAP = {
        "FMCG":["fmcg", "consumer", "food", "snack", "beverage", "dairy", "grocery", "cpg"],
        "Automotive & Transportation":["automotive", "car", "bike", "transport", "vehicle", "ev", "scooter", "auto", "mobility", "motor"],
        "Food & Beverage":["food", "beverage", "dairy", "snacks", "cafe", "restaurant", "fmcg", "dining", "qsr", "drink"],
        "Real Estate & Construction":["real estate", "builder", "property", "infra", "developer", "construction", "realty", "housing"],
        "Education & Training":["education", "school", "college", "edtech", "university", "learning", "institute", "academy"],
        "Healthcare":["health", "pharma", "fitness", "gym", "hospital", "wellness", "medical", "clinic", "care", "diagnostic"],
        "Pharma":["health", "pharma", "fitness", "gym", "hospital", "wellness", "medical", "clinic", "care", "medicine"],
        "Finance & Fintech":["finance", "bank", "insurance", "loan", "fintech", "wealth", "investment", "credit", "pay", "mutual fund"],
        "Retail":["retail", "fashion", "lifestyle", "store", "apparel", "luxury", "jewellery", "shop", "supermarket"],
        "E-Commerce":["ecommerce", "e-commerce", "online", "retail", "marketplace", "d2c", "delivery", "shopping"],
        "Technology & Business Services":["tech", "saas", "software", "it", "b2b", "consulting", "service", "app"],
        "Beauty & Personal Care":["beauty", "cosmetic", "skin", "hair", "personal care", "salon", "grooming", "makeup", "fragrance"],
        "Home Goods & Electronics":["interior", "furniture", "home", "decor", "paint", "furnishing", "appliances", "electronics", "tv", "smart"],
        "Hospitality & Travel":["travel", "hotel", "hospitality", "tourism", "flight", "booking", "holiday", "resort", "airline"],
        "Marketing, Advertising & Media":["marketing", "advertising", "media", "agency", "ott", "entertainment", "broadcast"],
        "Apparel & Fashion":["apparel", "fashion", "clothing", "wear", "shoes", "retail", "garment"],
        "Jewellery":["jewel", "gold", "diamond", "retail", "luxury", "accessory"],
        "Membership & Local Services":["service", "membership", "local", "salon", "spa", "subscription"],
        "Pets & Pet Services":["pet", "dog", "cat", "vet", "animal"],
        "Gaming": ["gaming", "esports", "games", "entertainment"],
        "Logistics & Warehousing":["logistics", "delivery", "warehouse", "supply", "b2b", "transport"],
        "Energy, Renewables & Mining":["energy", "solar", "power", "renewable", "electric"],
        "Manufacturing & Industrial":["manufacturing", "industrial", "factory", "b2b", "production"],
        "Quick Commerce":["quick commerce", "qcommerce", "delivery", "grocery", "fmcg", "ecommerce", "blinkit", "zepto", "instamart"],
        "OTT":["ott", "streaming", "media", "entertainment", "movie", "video", "content"]
    }
    
    strict_keywords = STRICT_INDUSTRY_MAP.get(target_brand_industry, [target_brand_industry.lower()])
    target_brand_clean = current_target_brand_name.lower().strip()
    
    sub_category_keywords = current_meeting_data.get('sub_category_keywords',[])

    print(f"    Searching Campaigns/Case Studies for STRICT MATCH ONLY: Brand='{target_brand_clean}' OR Industry='{strict_keywords[:2]}'")

    # --- 3. PROCESS FILES (Campaigns & Case Studies) ---
    all_files_in_folder = list_files_in_gdrive_folder(drive_service, NBH_GDRIVE_FOLDER_ID)
    
    for item in all_files_in_folder:
        fname = item.get('name', '')
        fid = item['id']
        mtype = item.get('mimeType', '')

        if FILE_NAME_NBH_PREVIOUS_MEETINGS_GSHEET.lower() in fname.lower(): continue

        def get_cached_content():
            if fid not in GDRIVE_FILE_CACHE:
                print(f"    📥 Downloading {fname} from Drive (First time this run)...")
                GDRIVE_FILE_CACHE[fid] = get_structured_gdrive_file_data(drive_service, sheets_service, fid, fname, mtype)
            else:
                print(f"    ⚡ Using cached data for {fname}...")
            return GDRIVE_FILE_CACHE[fid]

        if FILE_NAME_PHYSICAL_CAMPAIGNS_GSHEET.lower() in fname.lower():
            content = get_cached_content()
            extracted_rows = extract_strict_campaigns_and_case_studies(content, fname, target_brand_clean, strict_keywords, sub_category_keywords, target_cities, target_depts, email_to_geo_map)
            if extracted_rows: data_buckets["physical_campaigns"].extend(extracted_rows)

        elif FILE_NAME_DIGITAL_CAMPAIGNS_GSHEET.lower() in fname.lower():
            content = get_cached_content()
            extracted_rows = extract_strict_campaigns_and_case_studies(content, fname, target_brand_clean, strict_keywords, sub_category_keywords, target_cities, target_depts, email_to_geo_map)
            if extracted_rows: data_buckets["digital_campaigns"].extend(extracted_rows)

        elif FILE_NAME_LATEST_CASE_STUDIES_GSHEET.lower() in fname.lower():
             content = get_cached_content()
             # Case studies don't use department matching, so we pass empty set for depts
             extracted_rows = extract_strict_campaigns_and_case_studies(content, fname, target_brand_clean, strict_keywords, sub_category_keywords, target_cities, set(), email_to_geo_map)
             if extracted_rows: data_buckets["case_studies"].extend(extracted_rows)

    # --- 4. FINAL STRING ASSEMBLY (Formatting for Prompt) ---
    campaigns_str = "## NBH CAMPAIGN EXAMPLES (From Live Sheets)\n"
    
    campaigns_str += "### PHYSICAL CAMPAIGNS\n"
    if data_buckets["physical_campaigns"]: campaigns_str += "\n".join(data_buckets["physical_campaigns"])
    else: campaigns_str += "DATA_EMPTY: No physical campaign data available for this specific brand/industry."
        
    campaigns_str += "\n\n### DIGITAL CAMPAIGNS\n"
    if data_buckets["digital_campaigns"]: campaigns_str += "\n".join(data_buckets["digital_campaigns"])
    else: campaigns_str += "DATA_EMPTY: No digital campaign data available for this specific brand/industry."

    case_studies_str = "\n\n## NBH CASE STUDIES (From Consolidated Sheet)\n"
    case_studies_str += "### RELEVANT CASE STUDIES\n"
    if data_buckets["case_studies"]: case_studies_str += "\n".join(data_buckets["case_studies"])
    else: case_studies_str += "DATA_EMPTY: No case study data available for this specific brand/industry."

    final_llm_string = (
        f"{history_context_str}\n\n"
        f"{campaigns_str}\n"
        f"{case_studies_str}\n"
    )
    
    return {
        "llm_summary_string": final_llm_string,
        "is_overall_direct_follow_up": is_overall_direct_follow_up,
        "has_other_past_interactions": has_other_past_interactions,
        "condensed_past_meetings_for_alert": condensed_past_meetings_for_alert
    }
# --- Calendar Processing ---
def get_upcoming_meetings(calendar_service, calendar_id='primary', time_delta_hours=96): # Process meetings in next 3 days
    now_utc = datetime.datetime.utcnow()
    time_min_str = (now_utc - datetime.timedelta(hours=48)).isoformat() + 'Z'
    time_max_str = (now_utc + datetime.timedelta(hours=time_delta_hours)).isoformat() + 'Z'
    
    print(f'Getting events between {time_min_str} and {time_max_str}')
    try:
        events_result = calendar_service.events().list(
            calendarId=calendar_id, timeMin=time_min_str, timeMax=time_max_str,
            singleEvents=True, orderBy='startTime',
            # q='brand.vmeet@nobroker.in' # This might filter too early if brandvmeet is added as resource
        ).execute()
        events = events_result.get('items', [])
        
        # ✅ FILTER OUT TASK EVENTS (Skip events starting with task markers)
        TASK_PREFIXES = ['✅ TASK', '☑ TASK', 'TASK:', '[TASK]', 'Action Items:', 'bhargav', 'bhargav demo']
        filtered_events = []
        
        for event in events:
            title = event.get('summary', '').strip()
            is_task = any(title.startswith(prefix) for prefix in TASK_PREFIXES)
            
            if not is_task:
                filtered_events.append(event)
            else:
                print(f"  ⏭️  Skipping task event: '{title}'")
        
        print(f"  ✅ Filtered {len(events) - len(filtered_events)} task events, {len(filtered_events)} meetings remaining")
        return filtered_events
        
    except HttpError as error:
        print(f'An error occurred fetching events: {error}')
        return []

def load_processed_event_ids():
    # ... (same as before) ...
    if not os.path.exists(PROCESSED_EVENTS_FILE): return set()
    with open(PROCESSED_EVENTS_FILE, 'r') as f: return set(line.strip() for line in f)

def save_processed_event_id(event_id):
    # ... (same as before) ...
    with open(PROCESSED_EVENTS_FILE, 'a') as f: f.write(event_id + '\n')

EVENT_TAG_PROCESSED = "[NBH_BRIEF_AGENT_PROCESSED_V1]"
EVENT_TAG_ALERT_SENT = "[NBH_LEADERSHIP_ALERT_SENT]" 

# --- GLOBAL CACHE FOR GDRIVE FILES ---
# Prevents downloading massive sheets multiple times per run
GDRIVE_FILE_CACHE = {}

def is_event_already_tagged(event_description):
    desc = event_description or ""
    
    # 1. If it's fully processed, skip it safely
    if EVENT_TAG_PROCESSED in desc:
        return True
        
    # 2. Check for our NEW timestamped lock
    lock_match = re.search(r'\[NBH_PROCESSING_IN_PROGRESS_(\d+)\]', desc)
    if lock_match:
        lock_time = int(lock_match.group(1))
        current_time = int(time.time())
        # If lock is older than 5 minutes (300 seconds), the previous run timed out on GCP. Override it.
        if current_time - lock_time > 300:
            print("  ⚠️ Found STALE lock from a timed-out run. Overriding...")
            return False
        else:
            return True # Still locked by an active, healthy run
            
    # 3. FIX FOR YOUR CURRENT BUG: If the old string-based lock is found, override it!
    if "[NBH_PROCESSING_IN_PROGRESS]" in desc:
        print("  ⚠️ Found OLD stuck lock. Overriding...")
        return False

    return False

def tag_event_as_processing(calendar_service, event_id, calendar_id='primary'):
    """Locks the event with a timestamp so we can detect crashes later."""
    if not calendar_service: return
    try:
        event = calendar_service.events().get(calendarId=calendar_id, eventId=event_id).execute(num_retries=3)
        description = event.get('description', '')
        
        current_time = int(time.time())
        new_lock_tag = f"[NBH_PROCESSING_IN_PROGRESS_{current_time}]"
        
        # Clean up any old or stuck locks first
        clean_desc = re.sub(r'\[NBH_PROCESSING_IN_PROGRESS.*?\]', '', description).strip()
        
        new_description = f"{clean_desc}\n\n{new_lock_tag}"
        calendar_service.events().patch(
            calendarId=calendar_id, eventId=event_id, body={'description': new_description}
        ).execute(num_retries=3)
        print(f"  🔒 Locked event {event_id} (Timestamp: {current_time}).")
    except Exception as e:
        print(f"  ⚠️ Network error locking event {event_id}: {e}")

def tag_event_alert_sent(calendar_service, event_id, calendar_id='primary'):
    """Tags the calendar event immediately after a leadership alert is sent to prevent duplicates."""
    if not calendar_service: return
    try:
        event = calendar_service.events().get(calendarId=calendar_id, eventId=event_id).execute(num_retries=3)
        description = event.get('description', '')
        if EVENT_TAG_ALERT_SENT not in description:
            new_description = f"{description}\n\n{EVENT_TAG_ALERT_SENT}"
            updated_event_body = {'description': new_description}
            calendar_service.events().patch(
                calendarId=calendar_id, eventId=event_id, body=updated_event_body
            ).execute(num_retries=3)
            print(f"  Tagged event {event_id} with LEADERSHIP_ALERT_SENT flag.")
    except Exception as e:
        print(f"  ⚠️ Network error tagging alert sent for event {event_id}: {e}")

def tag_event_as_processed(calendar_service, event_id, calendar_id='primary'):
    if not calendar_service: return
    try:
        event = calendar_service.events().get(calendarId=calendar_id, eventId=event_id).execute(num_retries=3)
        description = event.get('description', '')
        
        # Remove ANY processing lock (old text or new timestamped)
        clean_desc = re.sub(r'\[NBH_PROCESSING_IN_PROGRESS.*?\]', '', description).strip()
        
        if EVENT_TAG_PROCESSED not in clean_desc:
            new_description = f"{clean_desc}\n\n{EVENT_TAG_PROCESSED}"
            calendar_service.events().patch(
                calendarId=calendar_id, eventId=event_id, body={'description': new_description}
            ).execute(num_retries=3)
            print(f"  ✅ Tagged event {event_id} as fully processed in calendar.")
    except Exception as e:
        print(f"  ⚠️ Network error tagging event {event_id} (Ignored to prevent crash): {e}")

EVENT_REMINDER_SET_TAG = "[NBH_1HR_REMINDER_SET]"

def set_one_hour_email_reminder(calendar_service, event_id, calendar_id='primary'):
    if not calendar_service:
        print("  Calendar service not available to set reminder.")
        return False
    try:
        # ADDED num_retries=3 to auto-reconnect
        event = calendar_service.events().get(calendarId=calendar_id, eventId=event_id).execute(num_retries=3)
        
        description = event.get('description', '')
        if EVENT_REMINDER_SET_TAG in description:
            print(f"  1-hour email reminder already marked as set for event {event_id}.")
            return True

        reminders = event.get('reminders', {})
        overrides = reminders.get('overrides',[])
        
        has_one_hour_email_reminder = any(
            r.get('method') == 'email' and r.get('minutes') == 60 for r in overrides
        )

        if has_one_hour_email_reminder:
            print(f"  Event {event_id} already has a 60-minute email reminder.")
        else:
            print(f"  Adding 60-minute email reminder to event {event_id}.")
            overrides.append({'method': 'email', 'minutes': 60})
            body_update = {'reminders': {'useDefault': False, 'overrides': overrides}}
            calendar_service.events().patch(
                calendarId=calendar_id, eventId=event_id, body=body_update
            ).execute(num_retries=3)
            print(f"  Successfully added 60-minute email reminder for event {event_id}.")

        if EVENT_REMINDER_SET_TAG not in description:
            new_description = f"{description}\n{EVENT_REMINDER_SET_TAG}".strip()
            body_update_desc = {'description': new_description}
            calendar_service.events().patch(
                calendarId=calendar_id, eventId=event_id, body=body_update_desc
            ).execute(num_retries=3)
            print(f"  Tagged event {event_id} with {EVENT_REMINDER_SET_TAG}.")
        
        return True

    except HttpError as error:
        print(f"  An HTTP error occurred setting reminder for event {event_id}: {error}")
        return False
    except Exception as e:
        # SAFETY NET: Prevents crash if connection drops
        print(f"  ⚠️ Network error setting reminder for event {event_id} (Ignored to prevent crash): {e}")
        return False

def send_gmail_message(gmail_service, user_id, message_body):
    """Sends an email message using the Gmail API."""
    if not gmail_service:
        print("  Gmail service not available. Cannot send email.")
        return None
    try:
        # ADDED num_retries=3 to auto-reconnect
        message = (gmail_service.users().messages().send(userId=user_id, body=message_body).execute(num_retries=3))
        print(f'  Message Id: {message["id"]} sent.')
        return message
    except HttpError as error:
        print(f'  An HTTP error occurred sending email: {error}')
        return None
    except Exception as e:
        # SAFETY NET
        print(f'  ⚠️ Network error sending email (Ignored to prevent crash): {e}')
        return None



# --- Meeting Info Extraction ---

NBH_SERVICE_ACCOUNTS_TO_EXCLUDE = { # Emails to exclude from the displayed NBH attendee list
    AGENT_EMAIL.lower(),
    "pia.brand@nobroker.in",
    "pia@nobroker.in",
    "nbh.meeting@gmail.com",
    "meetings.regional@gmail.com" 
}


def extract_meeting_info(event, agent_email_global, nbh_service_accounts_to_exclude_global):
    """
    Extracts basic, non-inferential information from a calendar event.
    The brand name itself is NOT processed here; it is extracted by the LLM later.
    """
    event_id = event['id']
    summary = event.get('summary', 'No Title')
    start_str = event['start'].get('dateTime', event['start'].get('date'))
    start_time_obj = datetime.datetime.fromisoformat(start_str.replace('Z', '+00:00'))
    location = event.get('location', 'N/A')
    description = event.get('description', '')
    attendees = event.get('attendees', [])

    nbh_attendees = []
    brand_attendees_info = []

    is_agent_invited = any(attendee.get('email', '').lower() == agent_email_global.lower() for attendee in attendees)
    if not is_agent_invited:
        print(f"  Skipping event '{summary}': {agent_email_global} is not an attendee.")
        return None

    for attendee in attendees:
        email = attendee.get('email', '').lower()
        name = attendee.get('displayName', email.split('@')[0] if '@' in email else email)
        
        # Skip if this is an excluded service account (check BEFORE categorization)
        if email in nbh_service_accounts_to_exclude_global:
            continue
        
        # Categorize as NBH or Brand attendee
        if '@nobroker.in' in email:
            nbh_attendees.append({'email': email, 'name': name})
        elif email:
            brand_attendees_info.append({'name': name, 'email': email})
    # Removing this condition so that physical meetings do not get skipped
    # if not brand_attendees_info:
    #     print(f"  Skipping event '{summary}': No external attendees.")
    #     return "NO_EXTERNAL_ATTENDEES"

    return {
        'id': event_id,
        'title': summary, # Return the raw title
        'start_time_obj': start_time_obj,
        'start_time_str': start_time_obj.strftime("%Y-%m-%d %I:%M %p %Z (%A)"),
        'location': location,
        'description': description,
        'nbh_attendees': nbh_attendees,
        'brand_attendees_info': brand_attendees_info,
        'is_event_description_present_for_tagging': bool(event.get('description'))
    }


# --- Gemini LLM Integration ---
def configure_gemini():
    """
    Configures and returns a Gemini LLM model instance using the provided API key.
    
    Returns:
        A configured Gemini GenerativeModel object if successful, or None if configuration fails or the API key is missing.
    """
    if not GEMINI_API_KEY:
        print("GEMINI_API_KEY environment variable not set. LLM will not function.")
        return None
    try:
        client = genai.Client(api_key=GEMINI_API_KEY)
        # Using a specific model version. 1.5 Flash is faster and cheaper for many tasks.
        # For higher quality, consider 'gemini-1.5-pro-latest'.
        # model = genai.GenerativeModel('gemini-2.5-flash')
        print(f"Gemini model configured successfully.")
        return client
    except Exception as e:
        print(f"Error configuring Gemini API: {e}")
        return None  


def generate_brief_with_gemini(gemini_llm_client, YOUR_DETAILED_PROMPT_TEMPLATE_GEMINI, meeting_data, internal_data_summary_str):
    if not gemini_llm_client:
        return "Error: Gemini model not available."

    nbh_attendee_names_str = ", ".join([att['name'] for att in meeting_data['nbh_attendees']])
    brand_attendee_names_only_str = ", ".join([att['name'] for att in meeting_data['brand_attendees_info']])
    
    # Format attendees with LinkedIn & Recent Activity
    brand_attendees_with_linkedin_str = ""
    for att in meeting_data['brand_attendees_info']:
        linkedin_display = att.get('linkedin_url', '(LinkedIn Not Verified)')
        if linkedin_display and linkedin_display != '(LinkedIn Not Verified)':
            linkedin_display = f"[LinkedIn Profile]({linkedin_display})"
            
        post_str = ""
        post_url = att.get('recent_post_url')
        if post_url and str(post_url).lower() != 'none' and str(post_url).lower() != 'null':
            context = att.get('post_context', 'View Recent Brand Activity')
            if not context or str(context).lower() == 'none':
                context = 'View Recent Brand Activity'
            post_str = f" | 📢 **Recent Activity:** [{context}]({post_url})"
        
        brand_attendees_with_linkedin_str += f"- **{att['name']}** ({att['email']}) - {linkedin_display}{post_str}\n"
        
    brand_attendees_info_str = "; ".join([f"{att['name']} ({att['email']})" for att in meeting_data['brand_attendees_info']])

    # Format potential key contacts
    potential_contacts_str = ""
    key_contacts_list = meeting_data.get('potential_key_contacts', [])

    if key_contacts_list:
        potential_contacts_str = "**Found Key Contacts:**\n\n"
        for contact in key_contacts_list:
            linkedin_display = contact.get('linkedin_url', '(LinkedIn Not Verified)')
            if linkedin_display and linkedin_display != '(LinkedIn Not Verified)':
                linkedin_display = f"[LinkedIn Profile]({linkedin_display})"
            
            potential_contacts_str += f"- **{contact['name']}** - {contact['title']} - {linkedin_display}\n"
            potential_contacts_str += f"  - Why They Matter: {contact['reasoning']}\n\n"
    else:
        potential_contacts_str = "**No additional key contacts found through search.**\n\n"

    # Execute Serper search for recent campaign news across all production meetings
    print(f"    🔍 [Serper News] Running search for '{meeting_data['brand_name']}' India news...")
    campaign_search_data = execute_serper_search_api(f'"{meeting_data["brand_name"]}" (campaign OR launch OR marketing) India (2025 OR 2026)', num_results=3)
    enriched_internal_summary = (
        f"{internal_data_summary_str}\n\n"
        f"## RECENT BRAND PUBLIC NEWS & CAMPAIGNS (VERIFIED GOOGLE INDEX):\n"
        f"{campaign_search_data}\n"
    )

    prompt_filled = YOUR_DETAILED_PROMPT_TEMPLATE_GEMINI.format(
        MEETING_DATETIME=meeting_data['start_time_str'],
        MEETING_LOCATION=meeting_data['location'],
        BRAND_NAME=meeting_data['brand_name'],
        BRAND_ATTENDEES_NAMES=brand_attendee_names_only_str,
        NBH_ATTENDEES_NAMES=nbh_attendee_names_str,
        BRAND_NAME_FOR_BODY=meeting_data['brand_name'],
        MEETING_TITLE=meeting_data.get('title', 'N/A'),
        BRAND_ATTENDEES_FULL_DETAILS=brand_attendees_info_str,
        BRAND_ATTENDEES_WITH_LINKEDIN=brand_attendees_with_linkedin_str,
        POTENTIAL_KEY_CONTACTS=potential_contacts_str,
        INTERNAL_NBH_DATA_SUMMARY=enriched_internal_summary
    )
    
    # Tool-free config to save API costs and prevent grounding latency
    config = types.GenerateContentConfig(
        temperature=0.0,
        top_p=0.95,
        top_k=40,
        safety_settings=[
            types.SafetySetting(category=types.HarmCategory.HARM_CATEGORY_HARASSMENT, threshold=types.HarmBlockThreshold.BLOCK_MEDIUM_AND_ABOVE),
            types.SafetySetting(category=types.HarmCategory.HARM_CATEGORY_HATE_SPEECH, threshold=types.HarmBlockThreshold.BLOCK_MEDIUM_AND_ABOVE),
            types.SafetySetting(category=types.HarmCategory.HARM_CATEGORY_SEXUALLY_EXPLICIT, threshold=types.HarmBlockThreshold.BLOCK_MEDIUM_AND_ABOVE),
            types.SafetySetting(category=types.HarmCategory.HARM_CATEGORY_DANGEROUS_CONTENT, threshold=types.HarmBlockThreshold.BLOCK_MEDIUM_AND_ABOVE),
        ]
    )

    print(f"  Sending request to Gemini for brand: {meeting_data['brand_name']}...")
    try:
        response = gemini_llm_client.models.generate_content(
            model="gemini-2.5-flash",
            contents=prompt_filled,
            config=config
        )
        if response.prompt_feedback and response.prompt_feedback.block_reason:
            return f"Error: Prompt blocked by Gemini. Reason: {response.prompt_feedback.block_reason_message or response.prompt_feedback.block_reason}"

        if not response.candidates:
             return f"Error: Gemini returned no candidates."

        brief_content = ""
        for part in response.candidates[0].content.parts:
            brief_content += part.text
        
        if not brief_content.strip():
            return "Error: Gemini returned an empty brief."
            
        brief_content = re.sub(r'\[cite:.*?\]', '', brief_content)
        return brief_content
    except Exception as e:
        print(f"  Error during Gemini API call: {e}")
        return f"Error: Exception during Gemini call: {e}"
# =====================================================================
# UNIFIED HIGH-QUALITY IMAGE GENERATION WORKFLOW
# =====================================================================
def get_brand_visual_context(gemini_client, brand_name, industry, generated_brief=""):
    """
    Acts as the Creative Director: Analyzes the pre-meeting brief and industry context
    to extract active campaigns, select target audiences, and dynamically generate 
    the conversion-oriented caption for Panel 3.
    """
    if not gemini_client: 
        return None
    
    prompt = f"""
    You are an expert Brand Visual Strategist and Creative Director at a top ad agency. 
    Analyze the generated Pre-Meeting Brief below for the brand '{brand_name}' (Industry: {industry}):
    ---
    {generated_brief}
    ---
    
    Task:
    1. is_well_known: Set to true ONLY if '{brand_name}' is a widely recognized national or multinational brand with established guidelines, recognizable logos, and clear visual identifiers in India (e.g., McDonald's, KFC, Coca-Cola, Tanishq, Puma, Horlicks, Swiggy, Amazon). Set to false if the brand is highly localized, a minor regional outlet, or obscure.
    2. primary_colors: Identify their exact 2 primary brand colors (e.g., McDonald's is "Golden Yellow and Crimson Red", Tanishq is "Deep Maroon and Gold").
    3. THE PITCH / CREATIVE HOOK: Scan the provided brief carefully. Focus specifically on the real-world, active, and researched campaigns of the brand mentioned in the text. Avoid generic themes.
    4. CREATE THE VISUAL SCENE: Combine the identified creative campaign/theme with premium, minimalist visual staging. 
       - EXCLUSION: Do not depict physical booths, tents, kiosks, or sampling structures. Focus solely on clean static print ads and graphic layouts on the gate, inside the elevator, and on the mobile interface.
    5. CREATE A SLOGAN: Extract the exact slogan proposed in the brief's creative hook, or draft a short, impactful 2-to-3 word slogan that matches the strategic pitch.
    6. DEDUCE TARGET AUDIENCE (DYNAMIC & SEMANTIC MATCHING):
       - Create a concise description of 1 or 2 specific Indian residents standing near or interacting with the advertisement.
       - The characters must match the exact nature and purchase intent of the brand.
       - Match the industry naturally.
    7. GENERATE DYNAMIC PANEL 3 CAPTION (MUST BE UNDER 12 WORDS):
       - Ensure caption_panel3 (Conversion Stage) matches the industry conversion metric exactly:
         - FMCG / Food / Quick Commerce: "Bringing your brand back when residents are ready to order."
         - Education: "Prompting parent inquiries and campus tours at the point of decision."
         - Jewellery: "Driving luxury brand aspiration and store visits at the moment of intent."
         - Automotive: "Prompting test-drive bookings when residents evaluate their next vehicle."
         - Finance/Fintech: "Driving secure account sign-ups and premium applications at point of intent."
         - Home Goods/Electronics: "Inspiring home consultation bookings and retail visits during relaxed hours."
         - Real Estate: "Cultivating high-value site visits for prospective residential buyers."
         - Wellness/Fitness: "Encouraging trial class bookings as residents plan their wellness routines."
         - Hospitality/Travel: "Sparking leisure travel bookings and holiday planning in relaxed moments."
         - Pets: "Triggering pet care bookings and premium nutrition orders near home."
         - OTT/Entertainment: "Prompting subscription sign-ups when residents seek leisure entertainment."
         - Apparel/Fashion: "Inspiring seasonal wardrobe updates and direct e-commerce checkouts."
         - Others: Synthesize a highly strategic equivalent following this pattern.

    Return ONLY a valid JSON object:
    {{
        "is_well_known": true/false,
        "primary_colors": "...",
        "visual_scene": "...",
        "short_slogan": "...",
        "target_audience": "Describe the dynamic target consumer(s) in clean, photorealistic terms",
        "caption_panel3": "Caption text for Panel 3 under 12 words"
    }}
    """
    
    try:
        grounding_tool = types.Tool(google_search=types.GoogleSearch())
        config = types.GenerateContentConfig(temperature=0.2, tools=[grounding_tool]) 
        
        response = gemini_client.models.generate_content(model="gemini-2.5-flash", contents=prompt, config=config)
        result_text = response.text.strip()
        result_text = re.sub(r'```json\s*|\s*```', '', result_text).strip()
        return json.loads(result_text)
    except Exception as e:
        print(f"  Visual context extraction failed for {brand_name}: {e}")
        return None

from playwright.sync_api import sync_playwright

def generate_creative_with_gemini_image(gemini_client, brand_name, industry, visual_context):
    """
    Generates a professional 3-panel side-by-side marketing funnel mockup collage.
    - Panel 1 and Panel 2 retain strictly static, brand-agnostic marketing funnel captions.
    - Panel 3 applies dynamic captions specific to the brand's industry and conversion goal.
    """
    if not visual_context:
        print(f"  Skipping image generation for {brand_name}: Visual context was not extracted.")
        return None

    # Strict screening check to prevent generating hallucinated logos for obscure brands
    if not visual_context.get("is_well_known"):
        print(f"  Skipping image generation for {brand_name}: Brand designated as regional or obscure to avoid visual hallucination.")
        return None
        
    colors = visual_context.get("primary_colors", "vibrant colors")
    visual_scene = visual_context.get("visual_scene", "modern lifestyle imagery")
    short_slogan = visual_context.get("short_slogan", "Exclusive Offer")
    
    # Extract dynamic demographic values and fallbacks
    target_audience = visual_context.get(
        "target_audience", 
        "A modern Indian resident dressed in clean, smart-casual attire"
    )
    
    # STATIC CAPTIONS: Standardized funnel markers
    caption_p1 = "When residents enter, your brand is discovered."
    caption_p2 = "Capturing undivided attention during high-focus transit moments."
    
    # DYNAMIC CAPTION: Customized conversion-stage marker
    caption_p3 = visual_context.get(
        "caption_panel3", 
        "Bringing your brand back when residents are ready to order."
    )
    
    print(f"  🎯 Dynamic Target Audience: '{target_audience}'")
    print(f"  📝 Captions -> P1 (Static): '{caption_p1}' | P2 (Static): '{caption_p2}' | P3 (Dynamic): '{caption_p3}'")
    
    # EXACT DESIGN PROMPT WITH STATIC & DYNAMIC VALUES INJECTED
    image_prompt = f"""
    Create a highly professional, commercial-grade horizontal mockup collage with a clean, symmetric aspect ratio. 
    The layout consists of a clean, plain white top header banner, followed by exactly three vertical panels (columns) positioned side-by-side, separated by thin, clean, solid white lines.
    
    # TOP HEADER BANNER (Span across the entire width of the collage):
    - A clean, solid white top header banner is required at the very top of the collage.
    - Displays the text centered in a highly polished, clean modern sans-serif typeface in a dark charcoal gray color (#333333):
      "From Visibility to Conversion — Powered by NoBrokerHood"

    # OVERALL STYLE, COLOR & AESTHETIC:
    - Consistent professional color grading and editorial tone across all three panels.
    - Soft, subtle depth-of-field background blurs to focus viewer attention on ad placements.
    - Real-world physics: clean drop-shadows under the poster frames and panels to convey realistic depth.
    - The palette strictly reflects the brand's primary colors ({colors}).

    # PANEL 1 (LEFT COLUMN): OUTDOOR RESIDENTIAL GATEWAY (Awareness)
    - ENVIRONMENT: Street-level daylight view of a premium Indian apartment complex entrance with a standard black sliding iron residential gate. High-rise buildings are visible in the soft-focus background.
    - MEDIA SETUP: A standard-sized horizontal rectangular banner (proportional scale, approximately 4 feet wide by 2.5 feet tall) is mounted flat and cleanly centered horizontally on the black gate bars. It must look naturally scaled and realistic, leaving significant portions of the gate's black bars visible above, below, and on the sides (not oversized, and not covering the full height or width of the gate).
    - ARTWORK: Displays '{brand_name}' logo, the campaign scene ("{visual_scene}"), and slogan "{short_slogan}".
    - HUMAN INTERACTION: {target_audience} walking past the gate, caught in a natural, candid moment looking directly at the advertisement.
    - CAPTION: White text centered in a dark, semi-transparent horizontal strip at the bottom. The text must read exactly:
      "{caption_p1}"

    # PANEL 2 (MIDDLE COLUMN): CAPTIVE ELEVATOR CABIN (Recall)
    - ENVIRONMENT: Interior of a sleek passenger lift cabin with modern brushed silver steel walls.
    - MEDIA: A vertical poster mounted inside a thin, clean, minimalist aluminum snap frame with soft drop shadows on the wall panel.
    - ARTWORK: Vertical layout of '{brand_name}' campaign scene ("{visual_scene}"), logo, and slogan "{short_slogan}".
    - HUMAN INTERACTION: The exact same {target_audience} from Panel 1 standing naturally inside the lift, looking at the poster.
    - CAPTION: White text centered in a dark, semi-transparent horizontal strip at the bottom. The text must read exactly:
      "{caption_p2}"

    # PANEL 3 (RIGHT COLUMN): NATIVE MOBILE IN-APP PLACEMENT (Conversion)
    - ENVIRONMENT: Premium close-up photo of a modern smartphone held in a hand, against a softly blurred home background.
    - MEDIA: Smartphone screen displaying a dark-themed NoBrokerHood application interface.
    - UI DETAILS: A clean white delivery pre-approval card floats at the top with a green header. Directly below is a full-width vertical banner ad for '{brand_name}' showing the campaign scene ("{visual_scene}"), logo, and slogan "{short_slogan}".
    - CAPTION: White text centered in a dark, semi-transparent horizontal strip at the bottom. The text must read exactly:
      "{caption_p3}"
    """
    
    api_successful = False
    raw_data = None

    # =====================================================================
    # ATTEMPT 1: DIRECT API GENERATION (PRIMARY / DEFAULT PRODUCTION ENGINE)
    # =====================================================================
    print(f"  🎨 Generating image using primary Gemini API model 'gemini-3-pro-image-preview'...")
    try:
        response = gemini_client.models.generate_content(
            model="gemini-3-pro-image-preview",
            contents=image_prompt,
            config=types.GenerateContentConfig(
                response_modalities=["IMAGE", "TEXT"]
            )
        )

        if response.candidates:
            for part in response.candidates[0].content.parts:
                if part.inline_data and part.inline_data.mime_type.startswith("image/"):
                    raw_data = part.inline_data.data
                    
                    if isinstance(raw_data, str):
                        raw_data = base64.b64decode(raw_data)
                    elif isinstance(raw_data, bytes):
                        if not raw_data.startswith(b'\xff\xd8\xff'):
                            raw_data = base64.b64decode(raw_data)
                    
                    print(f"  ✅ Primary API image generation successful for {brand_name}.")
                    api_successful = True
                    return raw_data

    except Exception as api_err:
        print(f"  ⚠️ Primary API generation failed or limits exceeded: {api_err}. Transitioning to Playwright fallback...")

    # =====================================================================
    # ATTEMPT 2: GOOGLE FLOW WEB AUTOMATION (SECONDARY / BACKUP OPTION)
    # =====================================================================
    if not api_successful:
        login_url = os.getenv("FLOW_LOGIN_URL")
        username = os.getenv("FLOW_USERNAME")
        password = os.getenv("FLOW_PASSWORD")
        auth_state_file = "flow_auth_state.json"
        flow_prompt = f"Generate this visual setup for the brand '{brand_name}': {image_prompt}"

        if os.path.exists(auth_state_file):
            print(f"  🎬 Launching Playwright browser to generate mockup on Google Flow...")
            try:
                with sync_playwright() as p:
                    browser = p.chromium.launch(headless=True)
                    context = browser.new_context(storage_state=auth_state_file)
                    page = context.new_page()
                    
                    page.set_viewport_size({"width": 1280, "height": 800})
                    page.goto("https://labs.google/fx/tools/flow", wait_until="networkidle", timeout=45000)
                    
                    # Verify login state
                    if "signin" in page.url or page.locator("text=Choose an account").is_visible():
                        print("  ❌ Session expired. Please regenerate your 'flow_auth_state.json' file locally.")
                        browser.close()
                        raise Exception("Session expired.")
                    
                    print("    ➕ Creating a new workspace project...")
                    page.locator("text=New project").click()
                    page.wait_for_load_state("networkidle")
                    
                    if page.locator("text=Got it").is_visible():
                        page.locator("text=Got it").click()
                    
                    prompt_input_selector = "textarea[placeholder*='What do you want to create']"
                    page.wait_for_selector(prompt_input_selector, timeout=15000)
                    
                    print("    📝 Entering design instructions into the canvas input...")
                    page.fill(prompt_input_selector, flow_prompt)
                    page.press(prompt_input_selector, "Enter")
                    
                    print("    ⏳ Prompt submitted. Rendering canvas (up to 120s)...")
                    output_image_selector = "div[class*='media-card'] img, div[class*='image'] img"
                    page.wait_for_selector(output_image_selector, timeout=120000)
                    
                    img_srcs = page.locator("img").all_get_attributes("src")
                    browser.close()
                    
                    valid_src = None
                    for src in img_srcs:
                        if src and ("googleusercontent.com" in src or "labs.google" in src):
                            valid_src = src
                            break
                    
                    if valid_src:
                        img_res = requests.get(valid_src, timeout=15)
                        if img_res.status_code == 200:
                            raw_data = img_res.content
                            print(f"  ✅ Web UI image generated successfully via Google Flow!")
                            return raw_data
                            
            except Exception as automation_err:
                print(f"  ⚠️ Google Flow Web Automation failed: {automation_err}")
        else:
            print("  ⚠️ 'flow_auth_state.json' file missing. Bypassing automation fallback.")

    return None
# --- Email Sending ---
def create_email_message_with_image(sender, to_emails_list, subject, message_text_html, image_bytes=None, brand_name=None):
    """Creates an email message, assigning a clear, custom filename to attachments to prevent 'noname' display."""
    msg = EmailMessage()
    msg["To"] = ", ".join(to_emails_list)
    msg["From"] = sender
    msg["Subject"] = subject

    # Plain text fallback
    msg.set_content("Please view this email in an HTML-compatible client to see the full brief and creatives.")
    
    # Add the HTML version
    msg.add_alternative(message_text_html, subtype="html")

    # Attach and embed Mockup Image using add_related with custom filename
    if image_bytes:
        html_part = msg.get_body(preferencelist=("html",))
        if html_part:
            # Assign a dynamic filename instead of leaving it default/empty
            safe_brand_name = "".join([c if c.isalnum() else "_" for c in (brand_name or "Brand")])
            attachment_filename = f"{safe_brand_name}_X_NoBrokerHood.jpg"
            
            html_part.add_related(
                image_bytes,
                maintype="image",
                subtype="jpeg",
                cid="creative_image",
                filename=attachment_filename
            )

    raw_message = base64.urlsafe_b64encode(msg.as_bytes()).decode()
    return {'raw': raw_message}

def send_gmail_message(gmail_service, user_id, message_body):
    """Sends an email message using the Gmail API."""
    if not gmail_service:
        print("  Gmail service not available. Cannot send email.")
        return None
    try:
        message = (gmail_service.users().messages().send(userId=user_id, body=message_body).execute())
        print(f'  Message Id: {message["id"]} sent.')
        return message
    except HttpError as error:
        print(f'  An error occurred sending email: {error}')
        return None

def send_brief_email(gmail_service, meeting_data, brief_content, creative_image_bytes=None):
    """Sends the brief email, injecting the AI creative if available. Includes TEST MODE."""
    EXCLUDED_EMAILS = {AGENT_EMAIL.lower(), "pia.brand@nobroker.in", "pia.hood@nobroker.in", "meetings.regional@gmail.com"} 

    nbh_recipient_emails =[]
    attendees_list = meeting_data.get('nbh_attendees',[]) 
    if isinstance(attendees_list, list):
        for att in attendees_list:
            if isinstance(att, dict) and 'email' in att:
                attendee_email = att.get('email')
                if attendee_email and isinstance(attendee_email, str) and attendee_email.lower() not in EXCLUDED_EMAILS:
                    nbh_recipient_emails.append(attendee_email)
    
    # =====================================================================
    # TEST MODE LOGIC: Change "True" to "False" when ready to go live!
    # =====================================================================
    TEST_MODE = False # <-- TURNED OFF! Emails will now go to actual attendees.
    
    if TEST_MODE:
        print("  ⚠️ TEST MODE IS ON: Overriding recipients. Sending only to Admin.")
        nbh_recipient_emails =[ADMIN_EMAIL_FOR_NOTIFICATIONS]
    # =====================================================================

    if not nbh_recipient_emails:
        print(f"  No NBH recipients for '{meeting_data['title']}'. Brief not emailed.")
        return

    email_subject = f"[{'TEST' if TEST_MODE else 'Pre-Meeting Brief'}]: {meeting_data['title']} with {meeting_data['brand_name']}"
    
    # --- Build and Inject the Event ID Box with Bright Yellow styling ---
    event_id_val = meeting_data.get('id', 'N/A')
    event_id_box_html = (
        f'<div class="event-id-box">'
        f'Event ID: <span class="event-id-text">{event_id_val}</span>'
        f'</div>'
    )

    # --- NEW: Build the Top 100 Sites Orange Highlight Box ---
    top_sites_url = "https://docs.google.com/spreadsheets/d/1NiYih9q_Gb-D6lCUjd08eDrsVAAJFqo6SRkBk8vzgrI/edit?gid=0#gid=0"
    top_sites_box_html = (
        f'<div class="top-sites-box">'
        f'Sample list of Top 100 Sites: <a href="{top_sites_url}" class="top-sites-link" target="_blank">Click Here to View Sheet</a>'
        f'</div>'
    )

    # Stack both boxes cleanly
    combined_boxes_html = f"{event_id_box_html}<br>{top_sites_box_html}"

    # Search for the Brand Attendees line in the markdown and insert the combined HTML boxes directly beneath it
    brand_attendees_pattern = re.compile(r'(Brand Attendees\s*:.*?)(\n|$)', re.IGNORECASE)
    if brand_attendees_pattern.search(brief_content):
        modified_brief_content = brand_attendees_pattern.sub(rf'\1\n\n{combined_boxes_html}\n', brief_content)
    else:
        # Prepend to the top of the brief as a fallback if the pattern is not found
        modified_brief_content = f"{combined_boxes_html}\n\n{brief_content}"

    html_brief_content = markdown.markdown(modified_brief_content)
    
    # --- INJECT IMAGE HTML ---
    creative_html = ""
    if creative_image_bytes:
        print("  📸 Attaching image to email HTML...")
        creative_html = """
        <div style="margin-bottom: 30px; background-color: #f8fafc; border: 1px solid #e2e8f0; padding: 25px; border-radius: 8px;">
            <h3 style="color: #2b6cb0; font-family: 'Segoe UI', Roboto, Helvetica, Arial, sans-serif; font-size: 18px; font-weight: bold; text-transform: uppercase; margin-top: 0; border-bottom: 2px solid #bee3f8; padding-bottom: 10px; text-align: left;">
                💡 THOUGHT-STARTERS: CREATIVES BASED ON CURRENT MARKET INSIGHTS
            </h3>
            <p style="font-size: 14px; color: #4a5568; margin-bottom: 20px; font-style: italic; text-align: left;">
                A suggested 3-in-1 visual pitch (Gate, Lift, and App) based on the brand's current live campaigns and colors:
            </p>
            <center>
                <img src="cid:creative_image" style="max-width: 100%; height: auto; border-radius: 8px; box-shadow: 0 6px 12px rgba(0,0,0,0.15);" alt="Brand Creative">
            </center>
        </div>
        """

    email_body_html = f"""
    <html>
    <head>
    <style>
        body {{ font-family: 'Segoe UI', Roboto, Helvetica, Arial, sans-serif; color: #333333; line-height: 1.6; background-color: #f4f7f6; padding: 20px; margin: 0; }}
        .email-container {{ max-width: 800px; margin: 0 auto; background-color: #ffffff; padding: 35px; border: 1px solid #e2e8f0; border-radius: 8px; box-shadow: 0 4px 6px rgba(0,0,0,0.05); }}
        .main-title {{ color: #0066cc; font-size: 20px; font-weight: bold; text-transform: uppercase; border-bottom: 2px solid #0066cc; padding-bottom: 12px; margin-bottom: 25px; letter-spacing: 0.5px; }}
        .greeting {{ font-size: 15px; color: #2d3748; margin-bottom: 25px; }}
        h1, h2 {{ color: #1a365d; font-size: 16px; font-weight: bold; text-transform: uppercase; margin-top: 35px; margin-bottom: 15px; border-bottom: 1px solid #edf2f7; padding-bottom: 5px; }}
        h3 {{ color: #2b6cb0; font-size: 15px; font-weight: bold; margin-top: 20px; margin-bottom: 10px; }}
        p {{ font-size: 14px; margin-top: 0; margin-bottom: 15px; }}
        ul, ol {{ margin-top: 5px; margin-bottom: 15px; padding-left: 25px; }}
        li {{ font-size: 14px; margin-bottom: 8px; color: #4a5568; }}
        .footer {{ margin-top: 40px; padding-top: 20px; border-top: 1px solid #e2e8f0; font-size: 13px; color: #718096; }}
        .highlight-box {{ background-color: #f0f8ff; border: 1px solid #bee3f8; border-radius: 6px; padding: 15px 20px; margin-top: 30px; margin-bottom: 20px; }}
        
        /* Soft Google Yellow Event ID Box Styles */
        .event-id-box {{
            background-color: #fef7e0;
            border: 1px solid #fbbc04;
            color: #b06000;
            border-radius: 6px;
            padding: 10px 16px;
            margin-top: 15px;
            margin-bottom: 8px;
            display: inline-block;
            font-size: 14px;
            font-weight: bold;
            font-family: 'Segoe UI', Roboto, Helvetica, Arial, sans-serif;
        }}
        .event-id-text {{
            font-family: 'Courier New', Courier, monospace;
            font-weight: bold;
            background-color: #ffffff;
            padding: 2px 8px;
            border: 1px solid #fde293;
            border-radius: 4px;
            margin-left: 5px;
            color: #3c4043;
        }}

        /* NEW: Orange Highlight Box for Top 100 Sites */
        .top-sites-box {{
            background-color: #fff5eb;
            border: 1px solid #ff9800;
            color: #d84315;
            border-radius: 6px;
            padding: 10px 16px;
            margin-top: 5px;
            margin-bottom: 20px;
            display: inline-block;
            font-size: 14px;
            font-weight: bold;
            font-family: 'Segoe UI', Roboto, Helvetica, Arial, sans-serif;
        }}
        .top-sites-link {{
            color: #0066cc;
            text-decoration: underline;
            margin-left: 5px;
        }}
    </style>
    </head>
    <body>
        <div class="email-container">
            <div class="main-title">NBH PRE-MEETING BRIEF</div>
            
            <div class="greeting">
                <p>Hi Team,</p>
                <p>Please find the Pre-Meeting Brief and Intelligence Report for your upcoming meeting with <strong>{meeting_data['brand_name']}</strong>.</p>
            </div>
            
            <!-- IMAGE INJECTED HERE FIRST -->
            {creative_html}
            
            <!-- TEXT BRIEF INJECTED HERE SECOND -->
            <div class="markdown-content">
                {html_brief_content}
            </div>
            
            <div class="footer">
                <p>Best regards,<br><strong>NBH Meeting Prep Agent</strong></p>
            </div>
        </div>
    </body>
    </html>
    """
    
    email_message = create_email_message_with_image(
        sender=AGENT_EMAIL,
        to_emails_list=nbh_recipient_emails,
        subject=email_subject,
        message_text_html=email_body_html,
        image_bytes=creative_image_bytes,
        brand_name=meeting_data.get('brand_name') # Added parameter
    )
    print(f"  FINAL CHECK: Sending styled brief for '{meeting_data['title']}' TO: {nbh_recipient_emails}")
    send_gmail_message(gmail_service, 'me', email_message)

def send_notification_email(gmail_service, subject, body_html, recipient=ADMIN_EMAIL_FOR_NOTIFICATIONS):
    """Sends error/alert notifications using the new email builder."""
    if not recipient:
        print("  Admin notification email not set. Skipping notification.")
        return
    
    # Always send a copy to brandvmeet for record-keeping
    recipients = list(set([recipient, AGENT_EMAIL]))

    email_message = create_email_message_with_image(
        sender=AGENT_EMAIL,
        to_emails_list=recipients,
        subject=subject,
        message_text_html=body_html,
        image_bytes=None # No image for simple notifications
    )
    send_gmail_message(gmail_service, 'me', email_message)

# Here are some useful functions to update meeting data in the master sheet

def read_data_from_sheets(sheet_id, sheets_service, range):

    try:
        result = (
                sheets_service.spreadsheets()
                .values()
                .get(spreadsheetId=sheet_id, range=range)
                .execute()
            )
        sheet_data = result.get("values", [])
        print(f"{len(sheet_data)} rows retrieved")
        return sheet_data
    except HttpError as error:
        print(f"An error occurred: {error}")

# Then check which events have not been updated in the google sheets
def events_to_update(meeting_ids, events):
    events_to_update = []
    for event in events:
        arr = [event["id"]]
        if arr not in meeting_ids:
            events_to_update.append(event)
    if not events_to_update:
        print("No new meetings to update")
        return None
    else:
        return events_to_update
    

def update_events_in_sheets(sheet_id, events_to_update, sheets_service, excluded_emails, designations, email_to_geo_map, column_index_master):
    import re
    import time
    
    # Get Column indices for Master Sheet (Meeting_data)
    master_sheet_columns = read_data_from_sheets(sheet_id, sheets_service, "Meeting_data!A1:BZ1")[0]
    owner_col_master = column_index[f"{master_sheet_columns.index('Owner') + 1}"]
    owner_update_col_master = column_index[f"{master_sheet_columns.index('Owner sheet to be updated') + 1}"]
    main_part_col = column_index[f"{master_sheet_columns.index('Main participant') + 1}"]
    meeting_done_col = column_index[f"{master_sheet_columns.index('Meeting Done') + 1}"]

    # Get Column indices for Audit Sheet
    audit_sheet_columns = read_data_from_sheets(sheet_id, sheets_service, "Audit_and_Training!A1:BZ1")[0]
    owner_col_audit = column_index[f"{audit_sheet_columns.index('Owner') + 1}"]
    owner_update_col_audit = column_index[f"{audit_sheet_columns.index('Owner sheet to be updated') + 1}"]

    for i, event in enumerate(events_to_update):
        id = event["id"]
        title = event.get("summary", "Untitled Meeting")
        
        # --- NEW: Extract Organizer Email ---
        organizer_email = event.get("organizer", {}).get("email", "")
        
        date = event["start"].get("date", event["start"].get("dateTime"))
        if 'T' in date:
            date = datetime.datetime.fromisoformat(date).date().isoformat()
        
        attendees = event.get("attendees")
        nobroker_attendee = []
        client_attendee = []
        if attendees:
            emails = [attendee["email"] for attendee in attendees]
            for email in emails:
                if email.lower() in excluded_emails:
                    continue
                if "nobroker" in email:
                    nobroker_attendee.append(email)
                else:
                    client_attendee.append(email)
        
        row_data = [id, title, date, f"{nobroker_attendee}", f"{client_attendee}"]
        
        # --- CALCULATE OWNER ---
        owner = None
        for email in nobroker_attendee:
            owner, _ = get_sheet_owner_from_email(email)
            if owner: break
            
        # --- NEW: CALCULATE DEPT (Brand Size) AND CITY ---
        found_dept = ""
        found_city = ""
        for email in nobroker_attendee:
            clean_email = email.lower()
            if clean_email in email_to_geo_map:
                raw_dept = email_to_geo_map[clean_email].get('dept', '')
                
                # Clean Dept logic from your separate script
                if 'regional' in raw_dept.lower(): found_dept = 'Regional'
                elif 'national' in raw_dept.lower(): found_dept = 'National'
                elif 'execution' in raw_dept.lower(): continue
                else: found_dept = raw_dept
                
                found_city = email_to_geo_map[clean_email].get('city', '').title()
                break # Stop after finding the first valid match
        
        # --- CALCULATE MAIN PARTICIPANT ---
        main_participant = []
        dg = []
        for role in ['bm', 'rm', 'ch']:
            if not main_participant:
                for p in nobroker_attendee:
                    d = designations.get(p, None)
                    if d and d.lower() == role:
                        main_participant.append(p)
                        dg.append(d)
        
        participant_data = [f"{main_participant}", f"{dg}", "Not Conducted"]

        def append_and_update_tab(tab_name, is_audit=False):
            try:
                # 1. Append the core data (A:E) and get the exact row number back
                append_result = sheets_service.spreadsheets().values().append(
                    spreadsheetId=sheet_id,
                    range=f"{tab_name}!A:A",
                    valueInputOption="USER_ENTERED",
                    insertDataOption="INSERT_ROWS",
                    body={"values": [row_data]}
                ).execute()
                
                updated_range = append_result.get('updates', {}).get('updatedRange', '')
                if not updated_range: return
                
                match = re.search(r'!A(\d+)', updated_range)
                if not match: return
                exact_row = match.group(1)
                
                # 2. Update Owner and Participant on that EXACT row
                update_payloads = []
                
                if is_audit:
                    if owner:
                        update_payloads.append({"range": f"{tab_name}!{owner_col_audit}{exact_row}:{owner_update_col_audit}{exact_row}", "values": [[owner, "TRUE"]]})
                else: # Meeting_data
                    if owner:
                        update_payloads.append({"range": f"{tab_name}!{owner_col_master}{exact_row}:{owner_update_col_master}{exact_row}", "values": [[owner, "TRUE"]]})
                    update_payloads.append({"range": f"{tab_name}!{main_part_col}{exact_row}:{meeting_done_col}{exact_row}", "values": [participant_data]})
                    
                    # --- NEW: PUSH DEPT, CITY, AND ORGANIZER TO MASTER SHEET ---
                    col_brand_size = column_index_master.get('Brand Size', 'L')
                    col_city = column_index_master.get('City (Attendees)', 'BJ')
                    col_organizer = column_index_master.get('Organizer of the Meeting', 'BK')
                    
                    if found_dept:
                        update_payloads.append({"range": f"{tab_name}!{col_brand_size}{exact_row}:{col_brand_size}{exact_row}", "values": [[found_dept]]})
                    if found_city:
                        update_payloads.append({"range": f"{tab_name}!{col_city}{exact_row}:{col_city}{exact_row}", "values": [[found_city]]})
                    if organizer_email:
                        update_payloads.append({"range": f"{tab_name}!{col_organizer}{exact_row}:{col_organizer}{exact_row}", "values": [[organizer_email]]})

                # 3. Push the updates
                if update_payloads:
                    sheets_service.spreadsheets().values().batchUpdate(
                        spreadsheetId=sheet_id,
                        body={"valueInputOption": "USER_ENTERED", "data": update_payloads}
                    ).execute()
                    
            except HttpError as error:
                print(f"Error updating {tab_name} for {title}: {error}")

        # Execute for the single sheet's tabs
        append_and_update_tab("Meeting_data", is_audit=False)
        append_and_update_tab("Audit_and_Training", is_audit=True)

        print(f"Appended and updated row for: {title}")

        # Rate limit protection (Wait 65s every 10 meetings to prevent Google Quota bans)
        if (i+1) % 10 == 0:
            print("Rate limit protection: Sleep initiated for 65 seconds...")
            time.sleep(65)

# Function to create a google document for the brief
def create_google_doc_in_folder(drive_service, folder_id, doc_name):
    file_metadata = {
        'name': doc_name,
        'mimeType': 'application/vnd.google-apps.document',
        'parents': [folder_id]
    }
    created = drive_service.files().create(
        body=file_metadata,
        fields='id, name, parents'
    ).execute()
    print(f"Created Google Doc: {created['name']} (ID: {created['id']})")
    return created['id']

# Function to write content to a Google Doc
def write_into_doc(docs_service, doc_id, text):
    requests = [
        {
            'insertText': {
                'location': { 'index': 1 },
                'text': text
            }
        }
    ]
    
    try:
        docs_service.documents().batchUpdate(
            documentId=doc_id,
            body={'requests': requests}
        ).execute()
    except:
        print("An error occured while writing into google doc")


# =====================================================================
# PPT GENERATION HELPER FUNCTIONS REMOVED
# =====================================================================
# Clean slate: PPT helper logic deleted to avoid execution during Serper testing.


def get_sheet_owner_from_email(email):
    hcy = []
    if email in sheet_masters:
        owner = email
        hcy.append(owner)
        return owner, hcy
    if email in hierarchy:
        manager = hierarchy[email]
        owner, hcy = get_sheet_owner_from_email(manager)
        hcy.append(email)
        return owner, hcy
    if email not in hierarchy:
        return None, []

# --- Main Execution Logic ---
def main():
    """
    Main orchestration function for automated pre-meeting brief generation and notification.
    
    This function coordinates the end-to-end workflow for preparing and emailing pre-meeting briefs for upcoming client meetings managed by the agent account. It initializes required Google Workspace services and the Gemini LLM, fetches upcoming calendar events, and processes each event as follows:
    
    - Skips events already processed or tagged.
    - Extracts meeting details and identifies the brand and industry using the LLM.
    - Retrieves and summarizes relevant internal NBH data for the brand.
    - Determines if the meeting is a direct follow-up or involves separate historical threads, and sends leadership alert emails as needed.
    - Generates a detailed pre-meeting brief using the LLM and internal data.
    - Emails the brief to NBH attendees, tags the event as processed, sets a 1-hour reminder, and records the event as processed.
    
    Handles error conditions gracefully, including missing services, ambiguous brand extraction, and LLM failures, with appropriate notifications and fallback logic.
    """
    
    print(f"Script started at {datetime.datetime.now()}")
    print(f"Using NBH GDrive Folder ID: {NBH_GDRIVE_FOLDER_ID}")
    
    # Load environment variables
    master_sheet_id = "1wWwjvAwXCAnPH3cAXSCaXlyDA7E9h_YrP4jaNFP9qvY" # main sheet 
    calendar_token = os.getenv("CALENDAR_TOKEN")
    gmail_token = os.getenv("GMAIL_TOKEN")
    drive_token = os.getenv("DRIVE_TOKEN")
    sheets_token = os.getenv("SHEET_TOKEN")
    docs_token = os.getenv("DOCS_TOKEN")  # Token for Google Docs API

    # Initialize Google Services
    # Use a combined token file strategy or separate ones. Separate is fine.
    calendar_service = get_google_service('calendar', 'v3', SCOPES, calendar_token)
    gmail_service = get_google_service('gmail', 'v1', SCOPES, gmail_token)
    drive_service = get_google_service('drive', 'v3', SCOPES, drive_token)
    sheets_service = get_google_service('sheets', 'v4', SCOPES, sheets_token)
    docs_service = get_google_service('docs', 'v1', SCOPES, docs_token)  # Docs service for creating briefs
    gemini_llm_client = configure_gemini()

    # Fetching employees data
    hcy_sheet_id = '1HxJt35QHF8BB_I8HusPQuiCS5_IpkEm5zoOSu1kwkNw'
    hcy_data = read_data_from_sheets(hcy_sheet_id, sheets_service, "Sheet4!A:F")
    df_hcy = pd.DataFrame(hcy_data[1:], columns=hcy_data[0])

    # Constructing designations and Geo mapping dictionary
    designations = {}
    email_to_geo_map = {}

    for i, row in df_hcy.iterrows():
        employee = str(row.get("Official Email ID", "")).strip().lower()
        dg = row.get("Designation New", "")
        location_col = str(row.get("Location", "")).strip().lower() # e.g., 'pune-regional'
        
        if employee and employee != 'nan':
            designations[employee] = dg
            
            city = ""
            dept = ""
            
            # Extract City and Dept from format like "Pune-National" or "Mumbai-regional"
            if '-' in location_col:
                parts = [p.strip() for p in location_col.split('-')]
                city = parts[0]
                if len(parts) > 1:
                    dept = parts[1]
            else:
                city = location_col
                
            email_to_geo_map[employee] = {'city': city, 'dept': dept}

    # Fetching column headers for master sheet and audit sheet
    master_sheet_columns = read_data_from_sheets(master_sheet_id, sheets_service,  "Meeting_data!A1:BZ1")[0]  # Get the header row
    audit_sheet_columns = read_data_from_sheets(master_sheet_id, sheets_service, "Audit_and_Training!A1:BZ1")[0]
    # Create a mapping of column names to their 1-based index
    global column_index_master
    global column_index_audit
    column_index_master = {name: column_index[f"{i+1}"] for i, name in enumerate(master_sheet_columns)}
    column_index_audit = {name: column_index[f"{i+1}"] for i, name in enumerate(audit_sheet_columns)}

    prompts_sheet_id = "1_dKfSF_WkANgSNvFbMTR43By_sK74XKWUr9fTzire5s"
    pre_meeting_brief = "Pre_meeting_brief"
    rng = f"{pre_meeting_brief}!A2:A2"
    pre_meeting_brief_prompt = read_data_from_sheets(prompts_sheet_id, sheets_service, rng)
    
    # --- SAFETY NET FOR GOOGLE SHEETS API GLITCHES ---
    if pre_meeting_brief_prompt and len(pre_meeting_brief_prompt) > 0 and len(pre_meeting_brief_prompt[0]) > 0:
        YOUR_DETAILED_PROMPT_TEMPLATE_GEMINI = pre_meeting_brief_prompt[0][0]
    else:
        print("CRITICAL ERROR: Could not fetch the Prompt Template from Google Sheets (API Glitch).")
        print("Exiting safely. Cloud Scheduler will retry on the next run.")
        return # Exits the script gracefully instead of crashing


    if not calendar_service: # Critical service
        print("Exiting: Calendar service failed to initialize.")
        return

    upcoming_events = get_upcoming_meetings(calendar_service)
    if not upcoming_events:
        print('No upcoming events found for agent email that need processing.')
        return

    # Updating events in the master sheet

    meeting_ids = read_data_from_sheets(master_sheet_id, sheets_service, "Meeting_data!A2:A")

    events_to_update_list = events_to_update(meeting_ids, upcoming_events)

    if not events_to_update_list:
        print("No new meetings to update in master sheet.")
    else:
        print(f"{len(events_to_update_list)} new meetings found")
        # ADDED BATCH LIMIT TO PREVENT CLOUD RUN TIMEOUT
        events_to_update_list = events_to_update_list[:5] 
        print(f"Limiting to 5 Master Sheet updates this run to prevent timeouts.")
        update_events_in_sheets(master_sheet_id, events_to_update_list, sheets_service, NBH_SERVICE_ACCOUNTS_TO_EXCLUDE, designations, email_to_geo_map, column_index_master)
    

    updated_meeting_ids = read_data_from_sheets(master_sheet_id, sheets_service, "Meeting_data!A2:A")

    processed_ids_local_file = load_processed_event_ids()

    # ADDED BATCH LIMIT VARIABLES
    MAX_BRIEFS_PER_RUN = 3
    briefs_generated_this_run = 0

    for event_payload in upcoming_events:
        # CHECK BATCH LIMIT
        if briefs_generated_this_run >= MAX_BRIEFS_PER_RUN:
            print(f"\n⏸️ Reached limit of {MAX_BRIEFS_PER_RUN} briefs for this execution.")
            print("Stopping to prevent Cloud Run timeout. Will process the rest on the next trigger.")
            break
            
        event_id = event_payload['id']
        event_summary = event_payload.get('summary', 'No Title')
        event_description_for_tag_check = event_payload.get('description')

        print(f"\nProcessing event: '{event_summary}' (ID: {event_id})")

        # Step 1: Check if the event has already been processed or is currently being processed
        if is_event_already_tagged(event_description_for_tag_check):
            print(f"  Skipping event '{event_summary}': Already tagged as processed.")
            continue
        
        if event_id in processed_ids_local_file:
            print(f"  Skipping event '{event_summary}': Found in local processed file.")
            continue

        # 🚀 CRITICAL FIX: Lock the event immediately so concurrent Webhooks don't duplicate work
        tag_event_as_processing(calendar_service, event_id)
        # Step 2: Extract basic meeting info (attendees, raw title, etc.)
        meeting_data_result = extract_meeting_info(event_payload, AGENT_EMAIL, NBH_SERVICE_ACCOUNTS_TO_EXCLUDE)

        # Step 3: Handle the possible "skip" results from the extraction
        if meeting_data_result is None: # Case where agent is not an attendee
            print(f"  Skipping event '{event_summary}': Agent is not an attendee.")
            save_processed_event_id(event_id)
            tag_event_as_processed(calendar_service, event_id)
            continue
        # Skipping this condition so that physical meetings can also be processed
        # if meeting_data_result == "NO_EXTERNAL_ATTENDEES":
        #     print(f"  Event '{event_summary}': No external attendees. No brief needed.")
        #     save_processed_event_id(event_id)
        #     tag_event_as_processed(calendar_service, event_id)
        #     continue

        # Step 4: If we are here, extraction was successful. Assign the result to meeting_data.
        # This is the key fix: assign the dictionary before trying to use it.
        meeting_data = meeting_data_result

        # Step 5: Use the LLM to get the brand name and industry from the raw title
        print(f"  Using LLM to extract brand details from title: '{meeting_data['title']}'")
        
        # --- NEW SAFETY BLOCK FOR API QUOTA CRASHES ---
        try:
            brand_details = get_brand_details_from_title_with_llm(gemini_llm_client, meeting_data['title'])
            # Add a pause after this call too
            time.sleep(5) 
        except Exception as e:
            print(f"  CRITICAL API ERROR for '{meeting_data['title']}': {e}")
            if "429" in str(e) or "RESOURCE_EXHAUSTED" in str(e):
                print("  ⚠️ Quota Exceeded. Pausing briefly (5s) to let API recover...")
                time.sleep(5)
                # Skip this meeting, try the next one (it remains locked as processing, will be retried later if lock cleared manually, or we can just let it fail and remove lock)
                continue
            else:
                # If it's another error, try to continue with unknown brand
                brand_details = {"brand_name": "Unknown Brand", "industry": "Unknown"}
        # --- END SAFETY BLOCK ---

        # Step 6: Handle ambiguous result from the LLM
        if brand_details['brand_name'] == 'Unknown Brand' or brand_details['brand_name'].lower() == 'unknown':
            print(f"  Event '{meeting_data['title']}': Title is ambiguous for brand extraction by LLM.")
            # Your notification logic for ambiguous titles can go here if needed.
            # Example:
            # ambiguous_body_html = f"..."
            # send_notification_email(...)
            save_processed_event_id(event_id)
            tag_event_as_processed(calendar_service, event_id) # Tag it so we don't retry
            # Updating the unknown brand name and industry in the master sheet
            index_of_event = updated_meeting_ids.index([event_id]) + 2 # +2 because A1 is header and A2 is first data row
            print(f"  Updating master sheet for event ID '{event_id}' at row {index_of_event} with brand 'Unknown")
            update_values = [[brand_details['brand_name'], brand_details['industry']]]
            body = {
            "valueInputOption": 'USER_ENTERED',  # Use USER_ENTERED to allow date formatting
            "data":[
                {"range": f"Meeting_data!F{index_of_event}:G{index_of_event}", "values": update_values},
                {"range": f"Audit_and_Training!F{index_of_event}:G{index_of_event}", "values": update_values},
                ],
            }
            try:
                sheets_service.spreadsheets().values().batchUpdate(spreadsheetId=master_sheet_id, body=body).execute()
                print(f"  Master sheet updated successfully for event ID '{event_id}'.")
                
                print(f" Resetting flag to TRUE for updating owner's sheet for event ID '{event_id}'")
                values = [["TRUE"]]
                flag_body = {
                    "valueInputOption": 'USER_ENTERED',
                    "data":[
                        {"range": f"Meeting_data!{column_index_master['Owner sheet to be updated']}{index_of_event}:{column_index_master['Owner sheet to be updated']}{index_of_event}", "values": values},
                        {"range": f"Audit_and_Training!{column_index_audit['Owner sheet to be updated']}{index_of_event}:{column_index_audit['Owner sheet to be updated']}{index_of_event}", "values": values},
                    ],
                }
                sheets_service.spreadsheets().values().batchUpdate(spreadsheetId=master_sheet_id, body=flag_body).execute()
                print(f"  Owner sheet flag reset for '{event_id}'.")
                
            except HttpError as error:
                print(f"  Error updating master sheet for event ID '{event_id}': {error}")
            
            continue

        # Step 7: Merge the successful LLM results into the main meeting_data dictionary
        meeting_data.update(brand_details)

        # Log Agency vs Direct Brand Context
        if meeting_data.get('is_agency_meeting'):
            print(f"  🏢 AGENCY MEETING DETECTED ({meeting_data.get('agency_name')}). Pitching with portfolio brand: '{meeting_data['brand_name']}'")

        # ========== NEW CODE STARTS HERE ==========
        # Get LinkedIn profiles for brand attendees using Serper
        print(f"  📱 Fetching LinkedIn profiles for brand attendees...")
        brand_attendees_with_linkedin = get_brand_attendees_linkedin_info(
            meeting_data.get('brand_attendees_info', []),
            meeting_data['brand_name'],
            gemini_llm_client
        )
        
        # Replace the old brand attendees info with the new one that has LinkedIn URLs
        meeting_data['brand_attendees_info'] = brand_attendees_with_linkedin
        # ========== NEW CODE ENDS HERE ==========

        # ========== NEW CODE FOR KEY CONTACTS STARTS HERE ==========
        # Find potential key contacts (people NOT in the meeting) using Serper
        print(f"  🎯 Finding potential key contacts at {meeting_data['brand_name']}...")
        potential_key_contacts = find_potential_key_contacts(
            meeting_data['brand_name'],
            gemini_llm_client
        )
        
        # Add to meeting data
        meeting_data['potential_key_contacts'] = potential_key_contacts
        # ========== NEW CODE FOR KEY CONTACTS ENDS HERE ==========

        current_brand_name_for_meeting = meeting_data['brand_name']
        target_brand_industry = meeting_data['industry']

        # Updating the brand name and industry in the master sheet
        index_of_event = updated_meeting_ids.index([event_id]) + 2 # +2 because A1 is header and A2 is first data row
        print(f"  Updating master sheet for event ID '{event_id}' at row {index_of_event} with brand '{current_brand_name_for_meeting}' and industry '{target_brand_industry}'")
        update_values = [[current_brand_name_for_meeting, target_brand_industry]]
        body = {
            "valueInputOption": 'USER_ENTERED',  # Use USER_ENTERED to allow date formatting
            "data":[
                {"range": f"Meeting_data!F{index_of_event}:G{index_of_event}", "values": update_values},
                {"range": f"Audit_and_Training!F{index_of_event}:G{index_of_event}", "values": update_values},
                ],
            }
        try:
            sheets_service.spreadsheets().values().batchUpdate(spreadsheetId=master_sheet_id, body=body).execute()
            print(f"  Master sheet updated successfully for event ID '{event_id}'.")
            
            print(f" Resetting flag to TRUE for updating owner's sheet for event ID '{event_id}'")
            values = [["TRUE"]]
            flag_body = {
                "valueInputOption": 'USER_ENTERED',
                "data":[
                    {"range": f"Meeting_data!{column_index_master['Owner sheet to be updated']}{index_of_event}:{column_index_master['Owner sheet to be updated']}{index_of_event}", "values": values},
                    {"range": f"Audit_and_Training!{column_index_audit['Owner sheet to be updated']}{index_of_event}:{column_index_audit['Owner sheet to be updated']}{index_of_event}", "values": values},
                ],
            }
            sheets_service.spreadsheets().values().batchUpdate(spreadsheetId=master_sheet_id, body=flag_body).execute()
            print(f"  Owner sheet flag reset for event ID '{event_id}'.")
            
        except HttpError as error:
            print(f"  Error updating master sheet for event ID '{event_id}': {error}")

        
        print(f"  LLM identified Brand: '{current_brand_name_for_meeting}', Industry: '{target_brand_industry}'")

        # --- THIS IS THE CORRECTED AND SIMPLIFIED BLOCK ---
        
        # Step 1: Check if necessary services are available.
        if drive_service and sheets_service:
            # If services are available, call the function to get the real data.
            internal_data_result = get_internal_nbh_data_for_brand(
                drive_service=drive_service,
                sheets_service=sheets_service,
                gemini_llm_client=gemini_llm_client,
                current_target_brand_name=current_brand_name_for_meeting,
                target_brand_industry=target_brand_industry, 
                current_meeting_data=meeting_data,
                EXCLUDED_NBH_PSEUDO_NAMES_FOR_FOLLOWUP=EXCLUDED_NBH_PSEUDO_NAMES_FOR_FOLLOWUP,
                AGENT_EMAIL=AGENT_EMAIL,
                master_sheet_id=master_sheet_id,
                email_to_geo_map=email_to_geo_map # FIXED PARAMETER
            )
        else:
            # If services are NOT available, create the default/fallback structure.
            print(f"  Drive/Sheets service not available. Skipping internal data fetch for '{current_brand_name_for_meeting}'.")
            internal_data_result = {
                "llm_summary_string": "Internal NBH Data: Not fetched due to Drive/Sheets service issues.",
                "is_overall_direct_follow_up": False,
                "has_previous_interactions": False,
                "condensed_past_meetings_for_alert": []
            }
        
        # Step 2: Extract the summary string for the LLM brief from the result (either real or default).
        internal_nbh_data_for_brand_str = internal_data_result["llm_summary_string"]
        
        # Step 3: Now, use the result for the leadership alert logic.
        has_prev_interactions_in_main = internal_data_result.get("has_previous_interactions", False)
        is_overall_follow_up_in_main = internal_data_result.get("is_overall_direct_follow_up", False)
        
        
        
                
        # --- >>> LEADERSHIP ALERT LOGIC (FINAL, CORRECTED VERSION) <<< ---

        # Step 1: Extract the flags and data we need from the internal data check.
        is_direct_follow_up = internal_data_result.get("is_overall_direct_follow_up", False)
        has_other_interactions = internal_data_result.get("has_other_past_interactions", False)
        condensed_meetings_for_alert = internal_data_result.get("condensed_past_meetings_for_alert", [])

        # Helper variables for the email body
        upcoming_meeting_title = meeting_data.get('title', 'N/A')
        upcoming_nbh_attendees_list = [att['name'] for att in meeting_data.get('nbh_attendees', [])]
        upcoming_nbh_attendees_str = ", ".join(upcoming_nbh_attendees_list) if upcoming_nbh_attendees_list else "N/A"

        # Check if we already sent an alert for this specific event to prevent spam loops
        already_alerted = EVENT_TAG_ALERT_SENT in (event_description_for_tag_check or "")

        # SCENARIO 1: "Hybrid" Engagement - A follow-up, but other separate threads also exist.
        if is_direct_follow_up and has_other_interactions:
            if already_alerted:
                print("DEBUG: HYBRID SCENARIO DETECTED, but alert was already sent previously. Skipping duplicate email.")
            else:
                print("DEBUG: HYBRID SCENARIO DETECTED. Sending nuanced leadership alert.")
                
                alert_subject = f"FYI: Complex Engagement with {current_brand_name_for_meeting} (Follow-up & Separate Threads)"
                
                alert_body_html = f"""
                <html><head><style> body {{ font-family: Arial, sans-serif; }} li {{ margin-bottom: 8px; }} </style></head>
                <body>
                    <p>Hello Leadership Team,</p>
                    <p>A new meeting has been scheduled with <b>{current_brand_name_for_meeting}</b>. This engagement is complex and requires coordination:</p>
                    <ul style="list-style-type:square;">
                        <li>It appears to be a <b>direct follow-up</b> to some recent discussions.</li>
                        <li>However, there are also <b>other, separate historical interactions</b> with this brand.</li>
                    </ul>
                    <p><b>Upcoming Meeting Details:</b></p>
                    <ul>
                        <li><b>Title:</b> {upcoming_meeting_title}</li>
                        <li><b>NBH Attendees:</b> {upcoming_nbh_attendees_str}</li>
                    </ul>
                    <p>This highlights a need for internal coordination. Context on the separate past interactions is below for awareness:</p>
                    <ul>
                """
                if condensed_meetings_for_alert:
                    for past_mtg in condensed_meetings_for_alert:
                        alert_body_html += f"<li><b>{past_mtg['date']}:</b> {past_mtg['discussion_summary']} (NBH Team: {past_mtg['nbh_team']})</li>"
                alert_body_html += "</ul><p>Best regards,<br>NBH Meeting Prep Agent</p></body></html>"
                
                # --- CORRECT EMAIL SENDING LOGIC ---
                if gmail_service and leadership_emails:
                    email_message = create_email_message_with_image(
                        sender=AGENT_EMAIL,
                        to_emails_list=leadership_emails,
                        subject=alert_subject,
                        message_text_html=alert_body_html
                    )
                    send_gmail_message(gmail_service, 'me', email_message)
                    print(f"    Leadership alert for HYBRID scenario with {current_brand_name_for_meeting} sent.")
                    # TAG IMMEDIATELY SO NEXT CRON JOB DOESN'T SPAM
                    tag_event_alert_sent(calendar_service, event_id)
                else:
                    print(f"    WARNING: Leadership alert for {current_brand_name_for_meeting} NOT sent (Gmail service or recipient list unavailable).")


        # SCENARIO 2: "Purely Separate" Engagement - Not a follow-up, but other past interactions exist.
        elif has_other_interactions and not is_direct_follow_up:
            if already_alerted:
                print("DEBUG: PURELY SEPARATE THREAD DETECTED, but alert was already sent previously. Skipping duplicate email.")
            else:
                print("DEBUG: PURELY SEPARATE THREAD DETECTED. Sending standard leadership alert.")
                
                alert_subject = f"FYI: New Meeting Scheduled with Existing Brand - {current_brand_name_for_meeting}"
                
                alert_body_html = f"""
                <html><head><style> body {{ font-family: Arial, sans-serif; }} li {{ margin-bottom: 8px; }} </style></head>
                <body>
                    <p>Hello Leadership Team,</p>
                    <p>A new meeting has been scheduled with <b>{current_brand_name_for_meeting}</b>. This meeting does <b>NOT</b> appear to be a direct follow-up to recent discussions.</p>
                    <p>This could indicate a new opportunity or a new NBH team engaging with the client.</p>
                    <p><b>Upcoming Meeting Details:</b></p>
                    <ul>
                        <li><b>Title:</b> {upcoming_meeting_title}</li>
                        <li><b>NBH Attendees:</b> {upcoming_nbh_attendees_str}</li>
                    </ul>
                    <p><b>Summary of Past Interactions (for context):</b></p>
                    <ul>
                """
                if condensed_meetings_for_alert:
                    for past_mtg in condensed_meetings_for_alert:
                        alert_body_html += f"<li><b>{past_mtg['date']}:</b> {past_mtg['discussion_summary']} (NBH Team: {past_mtg['nbh_team']})</li>"
                alert_body_html += "</ul><p>Best regards,<br>NBH Meeting Prep Agent</p></body></html>"
                
                # --- CORRECT EMAIL SENDING LOGIC ---
                if gmail_service and leadership_emails:
                    email_message = create_email_message_with_image(
                        sender=AGENT_EMAIL,
                        to_emails_list=leadership_emails,
                        subject=alert_subject,
                        message_text_html=alert_body_html
                    )
                    send_gmail_message(gmail_service, 'me', email_message)
                    print(f"    Leadership alert for SEPARATE THREAD with {current_brand_name_for_meeting} sent.")
                    # TAG IMMEDIATELY SO NEXT CRON JOB DOESN'T SPAM
                    tag_event_alert_sent(calendar_service, event_id)
                else:
                    print(f"    WARNING: Leadership alert for {current_brand_name_for_meeting} NOT sent (Gmail service or recipient list unavailable).")

        else:
            # This covers the "clean" cases: a brand-new meeting or a simple follow-up with no other threads.
            print("DEBUG: No leadership alert needed (Clean new meeting or simple follow-up).")
    


        if not gemini_llm_client:
            print(f"  Skipping brief generation for '{meeting_data['title']}': Gemini LLM not available.")
            # Don't mark as processed yet, maybe LLM will be available next run
            continue
        
        if not meeting_data.get('nbh_attendees'): # Check if any NBH humans are there
            print(f"  Event '{meeting_data['title']}': No NBH attendees (other than brandvmeet) to send brief to.")
            save_processed_event_id(event_id)
            tag_event_as_processed(calendar_service, event_id)
            continue


        print(f"  Proceeding with brief generation for: {meeting_data['brand_name']}")
        
        # 1. THE WRITER: Generate the Text Brief FIRST using Serper news context
        generated_brief = generate_brief_with_gemini(
            gemini_llm_client, 
            YOUR_DETAILED_PROMPT_TEMPLATE_GEMINI, 
            meeting_data, 
            internal_nbh_data_for_brand_str
        )

        # 2. IMAGE GENERATION (NOW LIVE FOR ALL MEETINGS)
        creative_image_bytes = None

        if ENABLE_IMAGE_GENERATION:
            if generated_brief and "Error:" not in generated_brief:
                try:
                    print(f"  🎨 Generating strategic mockup image for '{meeting_data['title']}'...")
                    
                    # Extract structured visual context from the text brief details (Art Director Role)
                    visual_context = get_brand_visual_context(
                        gemini_llm_client, 
                        meeting_data['brand_name'], 
                        meeting_data['industry'], 
                        generated_brief
                    )
                    
                    if visual_context:
                        # Render the standard 3-in-1 creative image layout for email inclusion
                        creative_image_bytes = generate_creative_with_gemini_image(
                            gemini_llm_client, 
                            meeting_data['brand_name'], 
                            meeting_data['industry'], 
                            visual_context
                        )
                        
                    else:
                        print(f"  Could not derive visual context for '{meeting_data['brand_name']}'.")
                
                except Exception as e:
                    print(f"  Warning: Failed to generate creative image: {e}")
            else:
                print(f"  ⚠️ Brief generation failed. Skipping image generation.")
        else:
            print(f"  ℹ️ Image Generation is disabled in settings. Skipping...")

        FEEDBACK_FORM_URL = "https://forms.gle/Ho9XLKsuGYhWBrBw7"

        # 2. Define the Footer Text (Using HTML injection inside Markdown)
        # We assign it the "highlight-box" class defined in our email CSS
        feedback_footer = f"""
\n\n
<div class="highlight-box">
    <p style="margin-bottom: 5px; font-size: 15px;"><strong>We want to hear from you!</strong></p>
    <p style="margin-bottom: 0;">Give your feedback on the Pre-Meeting Briefs.<br>
    <a href="{FEEDBACK_FORM_URL}">👉 Click Here to Fill the Feedback Form</a></p>
</div>
"""

        # 3. Append the footer to the generated brief
        # Only add it if the brief was generated successfully (no errors)
        if generated_brief and "Error:" not in generated_brief:
            generated_brief += feedback_footer

        if "Error:" in generated_brief or not generated_brief.strip(): # Check for errors from LLM
            print(f"  Failed to generate brief for '{meeting_data['title']}': {generated_brief}")
            error_body_html = f"""
            <html><body><p>The pre-meeting brief agent encountered an error while generating the brief for:</p>
            <p><b>Event:</b> {meeting_data['title']}<br>
            <b>Brand:</b> {meeting_data['brand_name']}<br>
            <b>Scheduled:</b> {meeting_data['start_time_str']}</p>
            <p><b>Error details:</b> {generated_brief}</p></body></html>"""
            send_notification_email(gmail_service,
                                    f"Error Generating Brief: {meeting_data['title']}",
                                    error_body_html)
            # Don't tag as fully processed if LLM fails, maybe it's temporary.
            # Or use a different tag like [NBH_BRIEF_AGENT_ERROR_V1]
        else:
            print(f"  Successfully generated brief for '{meeting_data['title']}'.")
            
            # Send the live email to the actual attendees!
            send_brief_email(gmail_service, meeting_data, generated_brief, creative_image_bytes)
            
            # ALWAYS tag the event as processed so the bot moves forward
            tag_event_as_processed(calendar_service, event_id) 
            set_one_hour_email_reminder(calendar_service, event_id) 
            save_processed_event_id(event_id)
            
            # --- Create Google Doc for the brief ---
            BRIEF_FOLDER_ID = "1RhhsFq5NGC2QtHPj8FQaU5BfhxJR5R6I"
            doc_id = create_google_doc_in_folder(
                drive_service,
                BRIEF_FOLDER_ID,
                f"Pre-Meeting Brief - {meeting_data['brand_name']} - {meeting_data['title']}"
            )
            if doc_id:
                write_into_doc(docs_service, doc_id=doc_id, text=generated_brief)
                # Updating doc link in master sheet
                index_of_event = updated_meeting_ids.index([event_id]) + 2 # +2 because A1 is header and A2 is first data row
                update_values = [[f"https://docs.google.com/document/d/{doc_id}"]]
                try:
                    body = {
                        "valueInputOption": 'USER_ENTERED',
                        "data":[
                            {"range": f"Meeting_data!H{index_of_event}:H{index_of_event}", "values": update_values},
                            {"range": f"Audit_and_Training!H{index_of_event}:H{index_of_event}", "values": update_values},
                        ],
                    }
                    sheets_service.spreadsheets().values().batchUpdate(spreadsheetId=master_sheet_id, body=body).execute()
                    print(f"  Master sheet updated with Google Doc link for event ID '{event_id}'.")
                    
                    print(f" Resetting flag to TRUE for updating owner's sheet for event ID '{event_id}'")
                    values = [["TRUE"]]
                    flag_body = {
                        "valueInputOption": 'USER_ENTERED',
                        "data":[
                            {"range": f"Meeting_data!{column_index_master['Owner sheet to be updated']}{index_of_event}:{column_index_master['Owner sheet to be updated']}{index_of_event}", "values": values},
                            {"range": f"Audit_and_Training!{column_index_audit['Owner sheet to be updated']}{index_of_event}:{column_index_audit['Owner sheet to be updated']}{index_of_event}", "values": values},
                        ],
                    }
                    sheets_service.spreadsheets().values().batchUpdate(spreadsheetId=master_sheet_id, body=flag_body).execute()
                    print(f"  Owner sheet flag reset for event ID '{event_id}'.")
                    
                except HttpError as error:
                    print(f"  Error updating master sheet with Google Doc link for event ID '{event_id}': {error}")
                # If we have an alternate sheet, update it too
                print(f"  Google Doc created and content written for '{meeting_data['title']}'.")
            
            # INCREMENT COUNTER AFTER SUCCESSFUL GENERATION
            briefs_generated_this_run += 1
            
        

    #print(f"Script finished at {datetime.datetime.now()}")

if __name__ == '__main__':
    main()
